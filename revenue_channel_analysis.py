#!/usr/bin/env python3
"""
2026年上半年营收渠道分析 — McKinsey 风格 PPT
包含堆叠柱状图（Pattern #38）展示渠道占比变化
"""

import os, zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════════════════════
#  Constants
# ═══════════════════════════════════════════════════════════
NAVY       = RGBColor(0x05, 0x1C, 0x2C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LINE_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)

ACCENT_BLUE   = RGBColor(0x00, 0x6B, 0xA6)

BODY_SIZE       = Pt(14)
SUB_HEADER_SIZE = Pt(18)
HEADER_SIZE     = Pt(28)
TITLE_SIZE      = Pt(22)

SW = Inches(13.333)
SH = Inches(7.5)
LM = Inches(0.8)
CONTENT_W = Inches(11.733)  # 13.333 - 0.8*2

# ═══════════════════════════════════════════════════════════
#  Helper Functions
# ═══════════════════════════════════════════════════════════

def _clean_shape(shape):
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)


def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)


def add_text(slide, left, top, width, height, text, font_size=Pt(14),
             font_name='Arial', font_color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font='KaiTi', anchor=MSO_ANCHOR.TOP,
             line_spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        p.space_before = line_spacing if i > 0 else Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size.pt * 1.35)
        for run in p.runs:
            set_ea_font(run, ea_font)
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _clean_shape(shape)
    return shape


def add_hline(slide, x, y, length, color=BLACK, thickness=Pt(0.5)):
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, x, y, length, h, color)


def add_oval(slide, x, y, letter, size=Inches(0.45), bg=NAVY, fg=WHITE):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = letter
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        set_ea_font(run, 'KaiTi')
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    _clean_shape(c)
    return c


def add_action_title(slide, text, title_size=Pt(22)):
    add_text(slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.9),
             text, font_size=title_size, font_color=BLACK, bold=True,
             font_name='Georgia', ea_font='KaiTi', anchor=MSO_ANCHOR.MIDDLE)
    add_hline(slide, Inches(0.8), Inches(1.05), Inches(11.7), color=BLACK, thickness=Pt(0.5))


def add_source(slide, text, y=Inches(7.05)):
    add_text(slide, Inches(0.8), y, Inches(11), Inches(0.3),
             text, font_size=Pt(9), font_color=MED_GRAY)


def add_page_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
             f"{num}/{total}", font_size=Pt(9), font_color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)


def full_cleanup(outpath):
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)


# ═══════════════════════════════════════════════════════════
#  Build Presentation
# ═══════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BL = prs.slide_layouts[6]  # Blank layout

TOTAL_SLIDES = 8

# ───────────────────────────────────────────────────────────
#  Slide 1: Cover
# ───────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BL)
add_rect(s1, 0, 0, SW, Inches(0.05), NAVY)
add_text(s1, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
         '2026年上半年营收渠道分析', font_size=Pt(44), font_name='Georgia',
         font_color=NAVY, bold=True, alignment=PP_ALIGN.LEFT)
add_text(s1, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6),
         '渠道结构演变与战略优化建议', font_size=Pt(24),
         font_color=DARK_GRAY)
add_text(s1, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
         '战略分析报告  |  2026年7月', font_size=BODY_SIZE,
         font_color=MED_GRAY)
add_hline(s1, Inches(1), Inches(6.8), Inches(3), NAVY, Pt(2))

# ───────────────────────────────────────────────────────────
#  Slide 2: Table of Contents
# ───────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BL)
add_action_title(s2, '目录')
toc_items = [
    ('1', '执行摘要', '核心发现与关键数字'),
    ('2', '渠道占比变化趋势', '1-6月堆叠柱状图分析'),
    ('3', '各渠道深度分析', '线上、线下、合作方三大渠道表现'),
    ('4', '关键驱动因素', '渠道变化背后的核心原因'),
    ('5', '战略建议', '下半年渠道优化方向'),
    ('6', '行动计划', '具体实施步骤与时间表'),
]
iy = Inches(1.6)
for num, title, desc in toc_items:
    add_oval(s2, LM, iy, num, size=Inches(0.45))
    add_text(s2, LM + Inches(0.65), iy, Inches(4.0), Inches(0.4),
             title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
    add_text(s2, Inches(5.5), iy + Inches(0.05), Inches(6.5), Inches(0.4),
             desc, font_size=BODY_SIZE, font_color=MED_GRAY)
    iy += Inches(0.65)
    add_hline(s2, LM, iy, CONTENT_W, LINE_GRAY)
    iy += Inches(0.25)
add_source(s2, 'Source: 内部销售数据汇总')
add_page_number(s2, 2, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 3: Executive Summary
# ───────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BL)
add_action_title(s3, '线上渠道占比在半年内从 35% 跃升至 62%，成为第一大收入来源')

# 核心结论框 — 增大高度到 1.2" 并缩小字号到 Pt(16) 防止文字溢出
add_rect(s3, LM, Inches(1.4), CONTENT_W, Inches(1.2), NAVY)
add_text(s3, LM + Inches(0.3), Inches(1.4), CONTENT_W - Inches(0.6), Inches(1.2),
         '核心发现：线上渠道在 6 个月内完成从第二大渠道到第一大渠道的逆转，占比从 35% 增长至 62%，线下门店从 45% 收缩至 23%',
         font_size=Pt(16), font_color=WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)

# 三大关键数字
stats = [
    ('+27pp', '线上渠道增长', '35%→62%', True),
    ('-22pp', '线下门店收缩', '45%→23%', False),
    ('-5pp',  '合作方微降', '20%→15%', False),
]
sw = Inches(3.5)
sg = (CONTENT_W - sw * 3) / 2
for i, (big, label, detail, is_navy) in enumerate(stats):
    sx = LM + (sw + sg) * i
    fill = NAVY if is_navy else BG_GRAY
    bc = WHITE if is_navy else NAVY
    sc = WHITE if is_navy else DARK_GRAY
    dc = WHITE if is_navy else MED_GRAY
    add_rect(s3, sx, Inches(3.0), sw, Inches(1.8), fill)
    add_text(s3, sx + Inches(0.2), Inches(3.1), sw - Inches(0.4), Inches(0.7),
             big, font_size=Pt(28), font_color=bc, bold=True,
             font_name='Georgia', alignment=PP_ALIGN.CENTER)
    add_text(s3, sx + Inches(0.2), Inches(3.8), sw - Inches(0.4), Inches(0.4),
             label, font_size=BODY_SIZE, font_color=sc, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s3, sx + Inches(0.2), Inches(4.25), sw - Inches(0.4), Inches(0.35),
             detail, font_size=Pt(12), font_color=dc,
             alignment=PP_ALIGN.CENTER)

# 底部要点
add_rect(s3, LM, Inches(5.3), CONTENT_W, Inches(1.5), BG_GRAY)
add_text(s3, LM + Inches(0.3), Inches(5.4), Inches(1.5), Inches(0.4),
         '关键启示', font_size=BODY_SIZE, font_color=NAVY, bold=True)
add_text(s3, LM + Inches(0.3), Inches(5.9), CONTENT_W - Inches(0.6), Inches(0.7),
         ['• 线上渠道已成为绝对主力，需持续加大投入以巩固优势',
          '• 线下门店需重新定位为体验中心，而非主要销售渠道',
          '• 合作方渠道虽占比下降，但绝对收入保持稳定，可作为补充渠道维护'],
         font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(8))

add_source(s3, 'Source: 财务部门月度报表，2026年1-6月')
add_page_number(s3, 3, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 4: Stacked Bar Chart — 渠道占比变化趋势 ★
# ───────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BL)
add_action_title(s4, '渠道占比变化趋势：线上渠道在 6 个月内从 35% 稳步增长至 62%')

# ── 数据定义 ──
periods = ['1月', '2月', '3月', '4月', '5月', '6月']
categories = ['线上', '线下', '合作方']
cat_colors = [NAVY, ACCENT_BLUE, LINE_GRAY]
data = [
    [35, 45, 20],   # 1月
    [40, 40, 20],   # 2月
    [48, 32, 20],   # 3月
    [52, 28, 20],   # 4月
    [58, 25, 17],   # 5月
    [62, 23, 15],   # 6月
]

# ── 图表标题（左对齐，在 action title 下方）──
chart_subtitle_y = Inches(1.2)
add_text(s4, LM, chart_subtitle_y, Inches(8.0), Inches(0.4),
         '各渠道收入占比月度变化（%）', font_size=Pt(13), font_color=DARK_GRAY,
         bold=True, alignment=PP_ALIGN.LEFT)

# ── 图例（右上角，紧凑排列，与图表右侧对齐）──
legend_y_top = chart_subtitle_y
chart_right_edge = LM + CONTENT_W  # 与页面内容区右边缘对齐
legend_right_edge = chart_right_edge
legend_item_spacing = Inches(1.4)  # 图例项之间的间距（紧凑）
legend_total_w = legend_item_spacing * len(categories) - Inches(0.2)
legend_start_x = legend_right_edge - legend_total_w
for ci, cat in enumerate(categories):
    lx = legend_start_x + legend_item_spacing * ci
    add_rect(s4, lx, legend_y_top + Inches(0.1), Inches(0.18), Inches(0.18), cat_colors[ci])
    add_text(s4, lx + Inches(0.25), legend_y_top, Inches(1.0), Inches(0.35),
             cat, font_size=Pt(11), font_color=DARK_GRAY)

# ── 图表区域参数 ──
chart_left = LM + Inches(0.8)
chart_right = chart_right_edge  # 与内容区右边缘对齐，消除右侧空白
chart_width = chart_right - chart_left

# ── 图表区域参数（整体下移，拉长图表填满页面）──
chart_top = Inches(2.0)
chart_bottom = Inches(5.7)
chart_height = chart_bottom - chart_top

n_periods = len(periods)
bar_width = Inches(0.72)  # 柱宽缩小40%（1.2 * 0.6 = 0.72）
bar_spacing = chart_width / n_periods

max_val = 100

# ── Y 轴刻度标签 + 参考线 ──
y_ticks = [0, 25, 50, 75, 100]
for tick in y_ticks:
    tick_y = chart_bottom - chart_height * (tick / max_val)
    add_text(s4, LM, tick_y - Inches(0.15), Inches(0.7), Inches(0.3),
             f'{tick}%', font_size=Pt(9), font_color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)
    if tick > 0:
        add_hline(s4, chart_left, tick_y, chart_width, LINE_GRAY, Pt(0.25))

# ── X 轴基线 ──
add_hline(s4, chart_left, chart_bottom, chart_width, BLACK, Pt(0.5))

# ── 绘制堆叠柱子 ──
for pi, period in enumerate(periods):
    bar_x = chart_left + bar_spacing * pi + (bar_spacing - bar_width) / 2
    cumulative = 0
    for ci in range(len(categories)):
        val = data[pi][ci]
        seg_h = chart_height * (val / max_val)
        seg_y = chart_bottom - chart_height * ((cumulative + val) / max_val)
        if val > 0:
            add_rect(s4, bar_x, seg_y, bar_width, seg_h, cat_colors[ci])
            # 段内百分比标签（当段高 >= 0.4" 时显示）
            if seg_h >= Inches(0.4):
                label_color = WHITE if ci <= 1 else DARK_GRAY
                add_text(s4, bar_x, seg_y, bar_width, seg_h,
                         f'{val}%', font_size=Pt(11), font_color=label_color,
                         bold=True, alignment=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)
        cumulative += val
    # X 轴标签
    add_text(s4, chart_left + bar_spacing * pi, chart_bottom + Inches(0.05),
             bar_spacing, Inches(0.3), period,
             font_size=BODY_SIZE, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ── 底部关键发现 ──
add_rect(s4, LM, Inches(6.15), CONTENT_W, Inches(0.75), BG_GRAY)
add_text(s4, LM + Inches(0.3), Inches(6.15), Inches(1.5), Inches(0.75),
         '关键发现', font_size=BODY_SIZE, font_color=NAVY, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s4, LM + Inches(2.0), Inches(6.15), CONTENT_W - Inches(2.3), Inches(0.75),
         '线上渠道占比从 1 月的 35% 稳步增长至 6 月的 62%，半年提升 27 个百分点；线下门店同期从 45% 收缩至 23%，渠道格局发生根本性逆转',
         font_size=BODY_SIZE, font_color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

add_source(s4, 'Source: 财务部门月度报表，2026年1-6月')
add_page_number(s4, 4, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 5: 各渠道深度分析（三支柱框架 #14）
# ───────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BL)
add_action_title(s5, '三大渠道各有不同发展态势，需差异化策略应对')

LIGHT_NAVY = RGBColor(0xE8, 0xEB, 0xEF)   # NAVY 的浅色版背景
LIGHT_BLUE2 = RGBColor(0xE3, 0xF0, 0xF7)  # ACCENT_BLUE 的浅色版背景
LIGHT_GRAY2 = RGBColor(0xF0, 0xF0, 0xF0)  # 灰色的浅色版背景

pillars = [
    ('线上渠道', NAVY, LIGHT_NAVY,
     ['占比：35% → 62%（+27pp）',
      '月均增长率：约 10%',
      '核心驱动：直播带货 + 社交裂变',
      '用户画像：25-35岁为主力',
      '客单价：¥280（同比+15%）']),
    ('线下门店', ACCENT_BLUE, LIGHT_BLUE2,
     ['占比：45% → 23%（-22pp）',
      '门店数量：从 120 家精简至 85 家',
      '坪效：同比下降 18%',
      '客流量：月均下降 12%',
      '优势：高客单价（¥520）']),
    ('合作方渠道', MED_GRAY, LIGHT_GRAY2,
     ['占比：20% → 15%（-5pp）',
      '合作商数量：维持 35 家',
      '绝对收入：基本持平',
      '账期：平均 45 天',
      '利润率：低于自营 8pp']),
]

pw = Inches(3.5)
pg = (CONTENT_W - pw * 3) / 2
for i, (title, accent, light_bg, points) in enumerate(pillars):
    px = LM + (pw + pg) * i
    # 顶部颜色条
    add_rect(s5, px, Inches(1.5), pw, Inches(0.06), accent)
    # 卡片背景
    add_rect(s5, px, Inches(1.56), pw, Inches(4.8), light_bg)
    # 标题
    add_text(s5, px + Inches(0.2), Inches(1.7), pw - Inches(0.4), Inches(0.5),
             title, font_size=SUB_HEADER_SIZE, font_color=accent, bold=True)
    add_hline(s5, px + Inches(0.2), Inches(2.3), pw - Inches(0.4), LINE_GRAY)
    # 内容
    add_text(s5, px + Inches(0.2), Inches(2.5), pw - Inches(0.4), Inches(3.5),
             [f'• {p}' for p in points],
             font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(8))

add_source(s5, 'Source: 各事业部月度经营报告，2026年H1')
add_page_number(s5, 5, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 6: 关键驱动因素（垂直步骤 #30）
# ───────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BL)
add_action_title(s6, '四大因素驱动渠道结构变化：数字化转型是核心推动力')

drivers = [
    ('1', '消费者行为迁移', '后疫情时代线上购物习惯固化，25-35岁消费者中 78% 优先选择线上渠道，移动端占比超过 85%'),
    ('2', '数字化营销投入', '上半年数字营销预算同比增加 60%，直播电商 GMV 增长 120%，社交媒体获客成本降低 35%'),
    ('3', '线下成本压力', '租金同比上涨 12%，人工成本增加 8%，单店盈亏平衡点提高 15%，倒逼门店精简优化'),
    ('4', '合作方政策调整', '收紧合作商准入门槛，淘汰低效合作方 5 家，聚焦头部合作商提升单产'),
]

iy = Inches(1.5)
for num, title, desc in drivers:
    add_oval(s6, LM, iy + Inches(0.05), num)
    add_text(s6, LM + Inches(0.65), iy, Inches(3.0), Inches(0.45),
             title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
    add_text(s6, Inches(4.8), iy, Inches(7.7), Inches(0.9),
             desc, font_size=BODY_SIZE, font_color=DARK_GRAY)
    iy += Inches(1.0)
    add_hline(s6, LM, iy, CONTENT_W, LINE_GRAY)
    iy += Inches(0.35)

# 底部总结
add_rect(s6, LM, Inches(5.9), CONTENT_W, Inches(0.9), BG_GRAY)
add_text(s6, LM + Inches(0.3), Inches(5.9), Inches(1.5), Inches(0.9),
         '趋势判断', font_size=BODY_SIZE, font_color=NAVY, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s6, LM + Inches(2.0), Inches(5.9), CONTENT_W - Inches(2.3), Inches(0.9),
         '渠道结构迁移为不可逆趋势，预计到 2026 年底线上渠道占比将达到 70%，线下门店将进一步精简至 60 家左右',
         font_size=BODY_SIZE, font_color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

add_source(s6, 'Source: 战略部行业研究报告，2026年6月')
add_page_number(s6, 6, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 7: 战略建议（行动计划 #35）
# ───────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BL)
add_action_title(s7, '下半年三大战略方向：加速线上、重塑线下、优化合作')

actions = [
    ('加速线上渠道', '2026 Q3-Q4',
     '投入 ¥2,000万 强化直播矩阵\n布局私域流量池（目标 50万会员）\n升级 AI 智能推荐引擎',
     '数字营销部'),
    ('重塑线下体验', '2026 Q3-Q4',
     '关闭 25 家低效门店\n核心门店升级为"体验+服务"中心\n线上线下会员体系打通',
     '零售运营部'),
    ('优化合作网络', '2026 H2',
     '聚焦 Top 15 核心合作商\n建立联合营销机制\n优化账期至 30 天',
     '渠道管理部'),
]

cw = Inches(3.5)
cg = (CONTENT_W - cw * 3) / 2
accents = [NAVY, ACCENT_BLUE, MED_GRAY]

for i, (title, timeline, desc, owner) in enumerate(actions):
    cx = LM + (cw + cg) * i
    # 顶部色条
    add_rect(s7, cx, Inches(1.5), cw, Inches(0.06), accents[i])
    # 标题栏
    add_rect(s7, cx, Inches(1.56), cw, Inches(0.6), NAVY)
    add_text(s7, cx + Inches(0.15), Inches(1.56), cw - Inches(0.3), Inches(0.6),
             title, font_size=BODY_SIZE, font_color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
    # 时间栏
    add_rect(s7, cx, Inches(2.16), cw, Inches(0.4), BG_GRAY)
    add_text(s7, cx + Inches(0.15), Inches(2.16), cw - Inches(0.3), Inches(0.4),
             timeline, font_size=BODY_SIZE, font_color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
    # 描述
    add_text(s7, cx + Inches(0.15), Inches(2.8), cw - Inches(0.3), Inches(2.5),
             desc.split('\n'), font_size=BODY_SIZE, font_color=DARK_GRAY,
             line_spacing=Pt(10), alignment=PP_ALIGN.LEFT)
    # 分隔线 + 负责人
    add_hline(s7, cx + Inches(0.3), Inches(5.5), cw - Inches(0.6), LINE_GRAY)
    add_text(s7, cx + Inches(0.15), Inches(5.7), cw - Inches(0.3), Inches(0.4),
             f'负责人：{owner}', font_size=BODY_SIZE, font_color=MED_GRAY,
             alignment=PP_ALIGN.CENTER)

add_source(s7, 'Source: 战略委员会决议，2026年7月')
add_page_number(s7, 7, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 8: Closing
# ───────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(BL)
add_rect(s8, 0, 0, SW, Inches(0.05), NAVY)
add_text(s8, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
         '拥抱渠道变革，赢得增长先机', font_size=Pt(28), font_color=NAVY, bold=True,
         font_name='Georgia', alignment=PP_ALIGN.CENTER)
add_hline(s8, Inches(5.5), Inches(3.3), Inches(2.3), NAVY, Pt(1.5))
add_text(s8, Inches(1.5), Inches(3.8), Inches(10.3), Inches(2.0),
         '线上渠道已成为未来增长的核心引擎\n以数据驱动决策，以体验重塑线下\n构建全渠道协同的新零售格局',
         font_size=SUB_HEADER_SIZE, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
         line_spacing=Pt(14))
add_hline(s8, Inches(1), Inches(6.8), Inches(3), NAVY, Pt(2))
add_page_number(s8, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════
#  Save & Cleanup
# ═══════════════════════════════════════════════════════════
outpath = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'output', '2026H1_营收渠道分析.pptx')
os.makedirs(os.path.dirname(outpath), exist_ok=True)
prs.save(outpath)
full_cleanup(outpath)
print(f'✅ Presentation saved to: {outpath}')
