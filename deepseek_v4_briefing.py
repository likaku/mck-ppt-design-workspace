#!/usr/bin/env python3
"""
DeepSeek-V4 发布动态周报 — McKinsey 风格 PPT
CONFIDENTIAL | 2026年3月4日 — 3月11日
"""

import os, zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════════════════════
#  Constants
# ═══════════════════════════════════════════════════════════
NAVY       = RGBColor(0x05, 0x1C, 0x2C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LINE_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)

ACCENT_BLUE   = RGBColor(0x00, 0x6B, 0xA6)

# 浅色背景版（四色方案对应的浅底色）
LIGHT_NAVY    = RGBColor(0xE8, 0xEB, 0xEF)   # NAVY 的浅色版
LIGHT_BLUE    = RGBColor(0xE3, 0xF0, 0xF7)   # ACCENT_BLUE 的浅色版
LIGHT_GRAY    = RGBColor(0xF0, 0xF0, 0xF0)   # MED_GRAY 的浅色版

BODY_SIZE       = Pt(14)
SUB_HEADER_SIZE = Pt(18)
HEADER_SIZE     = Pt(28)
TITLE_SIZE      = Pt(22)

SW = Inches(13.333)
SH = Inches(7.5)
LM = Inches(0.8)
CONTENT_W = Inches(11.733)  # 13.333 - 0.8*2

# ═══════════════════════════════════════════════════════════
#  Helper Functions
# ═══════════════════════════════════════════════════════════

def _clean_shape(shape):
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)


def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)


def add_text(slide, left, top, width, height, text, font_size=Pt(14),
             font_name='Arial', font_color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font='KaiTi', anchor=MSO_ANCHOR.TOP,
             line_spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        p.space_before = line_spacing if i > 0 else Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size.pt * 1.35)
        for run in p.runs:
            set_ea_font(run, ea_font)
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _clean_shape(shape)
    return shape


def add_hline(slide, x, y, length, color=BLACK, thickness=Pt(0.5)):
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, x, y, length, h, color)


def add_oval(slide, x, y, letter, size=Inches(0.45), bg=NAVY, fg=WHITE):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = letter
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        set_ea_font(run, 'KaiTi')
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    _clean_shape(c)
    return c


def add_action_title(slide, text, title_size=Pt(22)):
    add_text(slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.9),
             text, font_size=title_size, font_color=BLACK, bold=True,
             font_name='Georgia', ea_font='KaiTi', anchor=MSO_ANCHOR.MIDDLE)
    add_hline(slide, Inches(0.8), Inches(1.05), Inches(11.7), color=BLACK, thickness=Pt(0.5))


def add_source(slide, text, y=Inches(7.05)):
    add_text(slide, Inches(0.8), y, Inches(11), Inches(0.3),
             text, font_size=Pt(9), font_color=MED_GRAY)


def add_page_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
             f"{num}/{total}", font_size=Pt(9), font_color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)


def full_cleanup(outpath):
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)


# ═══════════════════════════════════════════════════════════
#  Build Presentation
# ═══════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BL = prs.slide_layouts[6]  # Blank layout

TOTAL_SLIDES = 13

# ───────────────────────────────────────────────────────────
#  Slide 1: Cover
# ───────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BL)
add_rect(s1, 0, 0, SW, Inches(0.05), NAVY)
add_text(s1, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
         'DeepSeek-V4 发布动态周报', font_size=Pt(44), font_name='Georgia',
         font_color=NAVY, bold=True, alignment=PP_ALIGN.LEFT)
add_text(s1, Inches(1), Inches(3.3), Inches(11.3), Inches(0.6),
         '从信息真空到全面临战——技术架构跃迁、算力主权化与资本市场联动',
         font_size=Pt(24), font_color=DARK_GRAY)
add_text(s1, Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
         'CONFIDENTIAL  |  2026年3月4日 — 3月11日', font_size=BODY_SIZE,
         font_color=MED_GRAY)
add_rect(s1, Inches(1), Inches(5.2), Inches(8.0), Inches(0.04), LINE_GRAY)
add_text(s1, Inches(1), Inches(5.5), Inches(11.3), Inches(0.8),
         ['技术架构跃迁：1M 上下文 + 原生多模态 + mHC/Engram 新架构',
          '算力主权化：战略性去英伟达化，华为昇腾深度绑定',
          '资本市场联动："DeepSeek 交易"成为全球半导体估值的宏观变量'],
         font_size=Pt(12), font_color=MED_GRAY, line_spacing=Pt(6))
add_hline(s1, Inches(1), Inches(6.8), Inches(3), NAVY, Pt(2))

# ───────────────────────────────────────────────────────────
#  Slide 2: Table of Contents
# ───────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BL)
add_action_title(s2, '目录')
toc_items = [
    ('1', '执行摘要', '核心结论与三维度影响评级矩阵'),
    ('2', '一周时间线', '从"信息真空"到"全面临战"的五个关键日期'),
    ('3', 'SCR 框架分析', 'Situation / Complication / Resolution 经典三段论'),
    ('4', '舆情可视化', '一周情绪变化趋势——分组柱状图'),
    ('5', '分角色行动建议', '按企业决策层/技术团队/投资者分别给出优先级'),
    ('6', '核心洞见', 'Key Takeaway 与战略意义'),
]
iy = Inches(1.6)
for num, title, desc in toc_items:
    add_oval(s2, LM, iy, num, size=Inches(0.45))
    add_text(s2, LM + Inches(0.65), iy, Inches(4.0), Inches(0.4),
             title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
    add_text(s2, Inches(5.5), iy + Inches(0.05), Inches(6.5), Inches(0.4),
             desc, font_size=BODY_SIZE, font_color=MED_GRAY)
    iy += Inches(0.65)
    add_hline(s2, LM, iy, CONTENT_W, LINE_GRAY)
    iy += Inches(0.25)
add_source(s2, 'Source: 知识库监控报告，2026-03-04 至 2026-03-11')
add_page_number(s2, 2, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 3: Executive Summary — 三维度矩阵
# ───────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BL)
add_action_title(s3, 'DeepSeek-V4 已进入发布倒计时，其战略意义已超越单纯的技术迭代')

# 核心结论框
add_rect(s3, LM, Inches(1.4), CONTENT_W, Inches(1.2), NAVY)
add_text(s3, LM + Inches(0.3), Inches(1.4), CONTENT_W - Inches(0.6), Inches(1.2),
         '核心结论：DeepSeek-V4 正在重新定义中国AI的"算力主权"路径，并成为影响全球半导体估值的宏观变量',
         font_size=Pt(16), font_color=WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)

# 三维度矩阵 — 数据表格
headers = ['维度', '核心变化', '影响评级']
col_widths = [Inches(2.5), Inches(6.5), Inches(2.733)]
hdr_y = Inches(2.9)
cx = LM
for hdr, cw in zip(headers, col_widths):
    add_text(s3, cx, hdr_y, cw, Inches(0.4), hdr,
             font_size=BODY_SIZE, font_color=MED_GRAY, bold=True)
    cx += cw
add_hline(s3, LM, Inches(3.35), CONTENT_W, BLACK, Pt(1.0))

rows = [
    ('技术架构', '1M上下文 + 原生多模态 + mHC/Engram新架构', '★★★★★'),
    ('算力生态', '战略性"去英伟达化"，华为昇腾优先适配', '★★★★★'),
    ('资本市场', '"DeepSeek交易"效应形成，全球半导体股价联动', '★★★★'),
]
row_h = Inches(0.7)
for ri, (dim, change, rating) in enumerate(rows):
    ry = Inches(3.5) + row_h * ri
    add_text(s3, LM, ry, col_widths[0], row_h, dim,
             font_size=BODY_SIZE, font_color=NAVY, bold=True)
    add_text(s3, LM + col_widths[0], ry, col_widths[1], row_h, change,
             font_size=BODY_SIZE, font_color=DARK_GRAY)
    add_text(s3, LM + col_widths[0] + col_widths[1], ry, col_widths[2], row_h, rating,
             font_size=BODY_SIZE, font_color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_hline(s3, LM, ry + row_h, CONTENT_W, LINE_GRAY)

# 底部三条主线
add_rect(s3, LM, Inches(5.8), CONTENT_W, Inches(1.0), BG_GRAY)
add_text(s3, LM + Inches(0.3), Inches(5.85), Inches(1.5), Inches(0.4),
         '三条主线', font_size=BODY_SIZE, font_color=NAVY, bold=True)
add_text(s3, LM + Inches(0.3), Inches(6.25), CONTENT_W - Inches(0.6), Inches(0.5),
         ['① 技术架构跃迁    ② 算力主权化    ③ 资本市场联动'],
         font_size=BODY_SIZE, font_color=DARK_GRAY)

add_source(s3, 'Source: 知识库监控报告（2026-03-04 至 2026-03-11），综合路透社、金融时报、科创板日报等')
add_page_number(s3, 3, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 4: 3月4日 — 信息真空期
# ───────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BL)
add_action_title(s4, '3月4日 — 信息真空期："暴风雨前的宁静"，市场期待感极度紧绷')

# 左侧：信号指标
add_rect(s4, LM, Inches(1.4), Inches(3.5), Inches(1.8), NAVY)
add_text(s4, LM + Inches(0.2), Inches(1.5), Inches(3.1), Inches(0.6),
         '248 条', font_size=Pt(28), font_color=WHITE, bold=True,
         font_name='Georgia', alignment=PP_ALIGN.CENTER)
add_text(s4, LM + Inches(0.2), Inches(2.2), Inches(3.1), Inches(0.8),
         '监控动态抓取量\n互动量全部为零（SEO空转）',
         font_size=Pt(12), font_color=WHITE, alignment=PP_ALIGN.CENTER)

# 右侧：关键事实
add_text(s4, Inches(5.0), Inches(1.5), Inches(7.5), Inches(0.4),
         '关键研判', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
add_hline(s4, Inches(5.0), Inches(2.0), Inches(7.5), LINE_GRAY)
add_text(s4, Inches(5.0), Inches(2.2), Inches(7.5), Inches(1.6),
         ['• 信号强度：低（噪音期），舆情中性（正面 20% / 中性 80%）',
          '• DeepSeek 处于极严信息保密期，与硅谷"渐进式泄露"完全不同',
          '• 高频关键词占位反映市场极度紧绷的期待感',
          '• 品牌溢价已形成——V4 发布将不仅是技术事件，而是商业分水岭'],
         font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(8))

# McKinsey Insight
add_rect(s4, LM, Inches(4.5), CONTENT_W, Inches(1.2), BG_GRAY)
add_text(s4, LM + Inches(0.3), Inches(4.55), Inches(2.0), Inches(0.4),
         'McKinsey Insight', font_size=BODY_SIZE, font_color=NAVY, bold=True)
add_text(s4, LM + Inches(0.3), Inches(5.0), CONTENT_W - Inches(0.6), Inches(0.6),
         '这种"真空态"在历史上是大型发布前的典型信号。积压的关注度一旦遇到真实信号，将引发指数级传播效应。',
         font_size=BODY_SIZE, font_color=DARK_GRAY)

add_source(s4, 'Source: 舆情监测平台，2026年3月4日')
add_page_number(s4, 4, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 5: 3月6日 — 信号爆发日
# ───────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BL)
add_action_title(s5, '3月6日 — 信号爆发日：三大重磅同日引爆，V4 从"模糊预期"变为"可定价事件"')

# 三支柱框架
pillars = [
    ('技术参数曝光', NAVY, LIGHT_NAVY,
     ['上下文窗口：1M Token',
      '（Claude 4.5 仅 200K）',
      'SWE-bench: 83.7%',
      '（Claude Opus 4.5: 80.9%）',
      '架构：MoE + Engram',
      '原生多模态（文/图/视频）']),
    ('"去英伟达化"公开', ACCENT_BLUE, LIGHT_BLUE,
     ['拒绝向英伟达/AMD',
      '提供V4早期访问权限',
      '优先适配：华为昇腾、寒武纪',
      '战略意图：构建"自主算法',
      '+ 自主算力"闭环']),
    ('"DeepSeek交易"升温', MED_GRAY, LIGHT_GRAY,
     ['A股概念股异动：',
      '浙江东方、润和软件、汉得信息',
      '华尔街担忧重演',
      '（R1曾致英伟达单日暴跌16%）',
      '舆情正面 75% / 负面 10%']),
]

pw = Inches(3.5)
pg = (CONTENT_W - pw * 3) / 2
for i, (title, accent, light_bg, points) in enumerate(pillars):
    px = LM + (pw + pg) * i
    add_rect(s5, px, Inches(1.3), pw, Inches(0.06), accent)
    add_rect(s5, px, Inches(1.36), pw, Inches(4.3), light_bg)
    add_text(s5, px + Inches(0.2), Inches(1.5), pw - Inches(0.4), Inches(0.5),
             title, font_size=SUB_HEADER_SIZE, font_color=accent, bold=True)
    add_hline(s5, px + Inches(0.2), Inches(2.1), pw - Inches(0.4), LINE_GRAY)
    add_text(s5, px + Inches(0.2), Inches(2.3), pw - Inches(0.4), Inches(3.0),
             [f'• {p}' for p in points],
             font_size=Pt(12), font_color=DARK_GRAY, line_spacing=Pt(4))

# McKinsey Insight
add_rect(s5, LM, Inches(5.9), CONTENT_W, Inches(0.9), BG_GRAY)
add_text(s5, LM + Inches(0.3), Inches(5.9), Inches(2.0), Inches(0.9),
         'McKinsey Insight', font_size=BODY_SIZE, font_color=NAVY, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s5, LM + Inches(2.5), Inches(5.9), CONTENT_W - Inches(2.8), Inches(0.9),
         '3月6日是整个追踪周期的"信息拐点"。从这一天起，V4 从"模糊预期"变为"可定价事件"。',
         font_size=BODY_SIZE, font_color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

add_source(s5, 'Source: 路透社、金融时报、科创板日报、微博、B站，2026年3月6日')
add_page_number(s5, 5, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 6: 3月9日 — 进入临战状态
# ───────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BL)
add_action_title(s6, '3月9日 — 进入临战状态：发布窗口锁定、架构革新、国产算力深化、Agent生态爆发')

# 四大关键发现 — 四栏
items = [
    ('1', '发布窗口锁定', ['路透社、金融时报多源确认', '万亿参数级旗舰', '同步发布技术说明']),
    ('2', '架构级革新', ['mHC：推理与知识分离', 'Engram：解决灾难性遗忘', '推理成本可降低达 90%']),
    ('3', '国产算力深化', ['昇腾推理速度提升 35x', 'Atlas 950超节点', '8192张昇腾950DT']),
    ('4', 'Agent生态爆发', ['OpenClaw 增速超 Linux', 'DeepSeek编码成功率极高', 'Token消耗量将指数增长']),
]
cw = Inches(2.7)
cg = (CONTENT_W - cw * 4) / 3
for i, (num, title, points) in enumerate(items):
    cx = LM + (cw + cg) * i
    add_rect(s6, cx, Inches(1.5), cw, Inches(4.3), BG_GRAY)
    add_oval(s6, cx + Inches(1.1), Inches(1.65), num)
    add_text(s6, cx + Inches(0.15), Inches(2.3), cw - Inches(0.3), Inches(0.4),
             title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_hline(s6, cx + Inches(0.3), Inches(2.8), cw - Inches(0.6), LINE_GRAY)
    add_text(s6, cx + Inches(0.15), Inches(3.0), cw - Inches(0.3), Inches(2.5),
             [f'• {p}' for p in points],
             font_size=Pt(12), font_color=DARK_GRAY, line_spacing=Pt(6))

# McKinsey Insight
add_rect(s6, LM, Inches(6.1), CONTENT_W, Inches(0.7), BG_GRAY)
add_text(s6, LM + Inches(0.3), Inches(6.1), Inches(2.0), Inches(0.7),
         'McKinsey Insight', font_size=BODY_SIZE, font_color=NAVY, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s6, LM + Inches(2.5), Inches(6.1), CONTENT_W - Inches(2.8), Inches(0.7),
         '竞争维度从"模型能力"升级到"模型+算力+Agent生态"三角闭环——从"单品发布"到"平台级事件"的质变。',
         font_size=BODY_SIZE, font_color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

add_source(s6, 'Source: 路透社、GitHub、华为昇腾技术文档，2026年3月9日')
add_page_number(s6, 6, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 7: 3月10日 — 技术确认日
# ───────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BL)
add_action_title(s7, '3月10日 — 技术确认日：灰度测试实锤 1M 上下文，产业链标的加速聚焦')

# 左栏：技术实锤
hw = Inches(5.5)
cg7 = Inches(0.733)
add_text(s7, LM, Inches(1.4), hw, Inches(0.4),
         '技术实锤', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
add_hline(s7, LM, Inches(1.9), hw, LINE_GRAY)
add_text(s7, LM, Inches(2.1), hw, Inches(3.0),
         ['✓  网页端/APP端已开始灰度测试 1M 上下文新模型',
          '✓  V4 Lite（Sealion-lite）效果显著优于现有模型',
          '✓  确认原生多模态（预训练阶段融合视觉+文本）',
          '',
          '产业链标的加速聚焦：',
          '• 高新发展、拓维信息（整机集成）',
          '• 华丰科技（连接器）、川润股份（液冷散热）',
          '• "华为昇腾+DeepSeek"被视为政企大模型标配'],
         font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(6))

# 右栏：风险信号
rx = LM + hw + cg7
add_text(s7, rx, Inches(1.4), hw, Inches(0.4),
         '风险信号', font_size=SUB_HEADER_SIZE, font_color=DARK_GRAY, bold=True)
add_hline(s7, rx, Inches(1.9), hw, DARK_GRAY)
add_text(s7, rx, Inches(2.1), hw, Inches(3.0),
         ['⚠  V4发布时间多次跳票',
          '    部分归因于GPU出口管控',
          '',
          '⚠  用户呼吁尽快上线API',
          '    灰度测试体验不佳',
          '',
          '⚠  "第三次跳票已经发生"',
          '    引发部分用户情绪反噬'],
         font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(6))

# McKinsey Insight
add_rect(s7, LM, Inches(5.5), CONTENT_W, Inches(0.9), BG_GRAY)
add_text(s7, LM + Inches(0.3), Inches(5.5), Inches(2.0), Inches(0.9),
         'McKinsey Insight', font_size=BODY_SIZE, font_color=NAVY, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s7, LM + Inches(2.5), Inches(5.5), CONTENT_W - Inches(2.8), Inches(0.9),
         '舆情正面 75% / 中性 20% / 负面 5%。信号强度极高（确认期）——灰度测试是从"预期"到"事实"的关键跨越。',
         font_size=BODY_SIZE, font_color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

add_source(s7, 'Source: DeepSeek 官方产品、A股公告、用户社区反馈，2026年3月10日')
add_page_number(s7, 7, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 8: 3月11日 — 临界点
# ───────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(BL)
add_action_title(s8, '3月11日 — 临界点：信号强度最高，全维度极度亢奋与焦虑并存')

# 三支柱：资本市场 / 开发者生态 / 风险信号
pillars8 = [
    ('资本市场', NAVY, LIGHT_NAVY,
     ['财经大V将V4视为3月最大催化剂',
      '"看DeepSeek V4硬度选股"成热词',
      'B站评测视频单条评论破千',
      '全网讨论达峰值']),
    ('开发者生态', ACCENT_BLUE, LIGHT_BLUE,
     ['长文本能力将对RAG创业',
      '产生降维打击',
      'AI Agent记忆能力获得质变',
      'OpenClaw本地部署成新热点']),
    ('风险信号', MED_GRAY, LIGHT_GRAY,
     ['"DeepSeek 又崩了"',
      '反映服务器稳定性隐忧',
      '"第三次跳票已经发生"',
      '引发部分用户情绪反噬']),
]

pw8 = Inches(3.5)
pg8 = (CONTENT_W - pw8 * 3) / 2
for i, (title, accent, light_bg, points) in enumerate(pillars8):
    px = LM + (pw8 + pg8) * i
    add_rect(s8, px, Inches(1.3), pw8, Inches(0.06), accent)
    add_rect(s8, px, Inches(1.36), pw8, Inches(3.8), light_bg)
    add_text(s8, px + Inches(0.2), Inches(1.5), pw8 - Inches(0.4), Inches(0.5),
             title, font_size=SUB_HEADER_SIZE, font_color=accent, bold=True)
    add_hline(s8, px + Inches(0.2), Inches(2.1), pw8 - Inches(0.4), LINE_GRAY)
    add_text(s8, px + Inches(0.2), Inches(2.3), pw8 - Inches(0.4), Inches(2.5),
             [f'• {p}' for p in points],
             font_size=Pt(12), font_color=DARK_GRAY, line_spacing=Pt(6))

# McKinsey Insight
add_rect(s8, LM, Inches(5.5), CONTENT_W, Inches(1.2), BG_GRAY)
add_text(s8, LM + Inches(0.3), Inches(5.55), Inches(2.0), Inches(0.4),
         'McKinsey Insight', font_size=BODY_SIZE, font_color=NAVY, bold=True)
add_text(s8, LM + Inches(0.3), Inches(6.0), CONTENT_W - Inches(0.6), Inches(0.6),
         '当市场情绪进入"极度亢奋"阶段，需警惕"利好出尽"风险。真正的价值判断应回归V4的技术验证结果与API定价策略。',
         font_size=BODY_SIZE, font_color=DARK_GRAY)

add_source(s8, 'Source: 微博、B站、投资社区、GitHub，2026年3月11日')
add_page_number(s8, 8, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 9: SCR 框架 — Case Study Pattern #34
# ───────────────────────────────────────────────────────────
s9 = prs.slides.add_slide(BL)
add_action_title(s9, 'SCR 框架分析：V4 技术领先但面临算力瓶颈与竞争窗口收窄双重挑战')

sections = [
    ('S', 'Situation\n现状', NAVY, LIGHT_NAVY,
     'V4 已完成从信息真空到全面临战的转变\n技术参数全面对标国际顶尖水平\n（1M上下文、原生多模态、SWE-bench 83.7%）'),
    ('C', 'Complication\n挑战', ACCENT_BLUE, LIGHT_BLUE,
     '算力瓶颈：GPU出口管控致训练延迟\n发布不确定性：多次跳票消耗市场耐心\n竞争窗口收窄：GPT-5、Claude 5 在即'),
    ('R', 'Resolution\n建议', MED_GRAY, LIGHT_GRAY,
     '保持技术栈灵活性，V4首日测试成本效益比\n重新评估RAG架构必要性\n关注真正完成V4适配的硬核标的'),
]

sw9 = Inches(3.5)
sg9 = (CONTENT_W - sw9 * 3) / 2
for i, (letter, title, accent, light_bg, desc) in enumerate(sections):
    sx = LM + (sw9 + sg9) * i
    add_rect(sx, sx, Inches(1.3), sw9, Inches(0.06), accent) if False else None
    add_rect(s9, sx, Inches(1.3), sw9, Inches(0.06), accent)
    add_rect(s9, sx, Inches(1.36), sw9, Inches(4.5), light_bg)
    add_oval(s9, sx + Inches(0.15), Inches(1.55), letter,
             bg=accent, fg=WHITE)
    add_text(s9, sx + Inches(0.15), Inches(2.1), sw9 - Inches(0.3), Inches(0.8),
             title, font_size=BODY_SIZE, font_color=accent, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_hline(s9, sx + Inches(0.2), Inches(3.0), sw9 - Inches(0.4), LINE_GRAY)
    add_text(s9, sx + Inches(0.2), Inches(3.2), sw9 - Inches(0.4), Inches(2.3),
             desc.split('\n'), font_size=Pt(12), font_color=DARK_GRAY,
             line_spacing=Pt(6))

# 底部总结
add_rect(s9, LM, Inches(6.1), CONTENT_W, Inches(0.7), BG_GRAY)
add_text(s9, LM + Inches(0.3), Inches(6.1), CONTENT_W - Inches(0.6), Inches(0.7),
         '关键判断：留给V4的"独占窗口期"有限，GPT-5、Claude 5 发布在即，速度即战略',
         font_size=BODY_SIZE, font_color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE, bold=True)

add_source(s9, 'Source: 综合分析，2026年3月4日—11日')
add_page_number(s9, 9, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 10: 舆情可视化 — 分组柱状图 (Pattern #37)
# ───────────────────────────────────────────────────────────
s10 = prs.slides.add_slide(BL)
add_action_title(s10, '一周舆情演变：情绪分布从3月4日的"冷静等待"急剧转向"全面看多"')

# 图表标题 + 图例
chart_subtitle_y = Inches(1.2)
add_text(s10, LM, chart_subtitle_y, Inches(6.0), Inches(0.4),
         '舆情情绪分布（%）', font_size=Pt(13), font_color=DARK_GRAY,
         bold=True, alignment=PP_ALIGN.LEFT)

# 数据定义
dates = ['3/4', '3/6', '3/9', '3/10', '3/11']
categories = ['正面', '中性', '负面']
cat_colors = [NAVY, LINE_GRAY, MED_GRAY]
data = [
    [20, 80, 0],    # 3/4
    [75, 15, 10],   # 3/6
    [75, 20, 5],    # 3/9
    [75, 20, 5],    # 3/10
    [75, 15, 10],   # 3/11
]

# 图例（右上角）
chart_right_edge = LM + CONTENT_W
legend_item_spacing = Inches(1.4)
legend_total_w = legend_item_spacing * len(categories) - Inches(0.2)
legend_start_x = chart_right_edge - legend_total_w
for ci, cat in enumerate(categories):
    lx = legend_start_x + legend_item_spacing * ci
    add_rect(s10, lx, chart_subtitle_y + Inches(0.1), Inches(0.18), Inches(0.18), cat_colors[ci])
    add_text(s10, lx + Inches(0.25), chart_subtitle_y, Inches(1.0), Inches(0.35),
             cat, font_size=Pt(11), font_color=DARK_GRAY)

# 图表区域参数
chart_left = LM + Inches(0.8)
chart_right = chart_right_edge
chart_width = chart_right - chart_left
chart_top = Inches(1.8)
chart_bottom = Inches(5.4)
chart_height = chart_bottom - chart_top

n_dates = len(dates)
n_cats = len(categories)
group_width = chart_width / n_dates
bar_width = Inches(0.4)
bar_gap = Inches(0.05)
group_bar_width = bar_width * n_cats + bar_gap * (n_cats - 1)
max_val = 100

# Y 轴刻度
y_ticks = [0, 20, 40, 60, 80, 100]
for tick in y_ticks:
    tick_y = chart_bottom - chart_height * (tick / max_val)
    add_text(s10, LM, tick_y - Inches(0.15), Inches(0.7), Inches(0.3),
             f'{tick}%', font_size=Pt(9), font_color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)
    if tick > 0:
        add_hline(s10, chart_left, tick_y, chart_width, LINE_GRAY, Pt(0.25))

# X 轴基线
add_hline(s10, chart_left, chart_bottom, chart_width, BLACK, Pt(0.5))

# 绘制柱子
for di, date in enumerate(dates):
    group_x = chart_left + group_width * di + (group_width - group_bar_width) / 2
    for ci, cat in enumerate(categories):
        val = data[di][ci]
        bar_h = chart_height * (val / max_val)
        bar_x = group_x + (bar_width + bar_gap) * ci
        bar_y = chart_bottom - bar_h
        if val > 0:
            add_rect(s10, bar_x, bar_y, bar_width, bar_h, cat_colors[ci])
            if val >= 10:
                add_text(s10, bar_x - Inches(0.05), bar_y - Inches(0.25),
                         bar_width + Inches(0.1), Inches(0.25),
                         f'{val}%', font_size=Pt(9), font_color=DARK_GRAY,
                         alignment=PP_ALIGN.CENTER)
    # X 轴日期标签
    add_text(s10, chart_left + group_width * di, chart_bottom + Inches(0.05),
             group_width, Inches(0.3), date,
             font_size=BODY_SIZE, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# 底部趋势总结
add_rect(s10, LM, Inches(5.9), CONTENT_W, Inches(0.8), BG_GRAY)
add_text(s10, LM + Inches(0.3), Inches(5.9), Inches(1.5), Inches(0.8),
         '趋势总结', font_size=BODY_SIZE, font_color=NAVY, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s10, LM + Inches(2.0), Inches(5.9), CONTENT_W - Inches(2.3), Inches(0.8),
         '舆情从3月4日的"冷静等待"急剧转向3月6日后的"全面看多"，负面情绪集中在发布延迟与服务稳定性两个点上',
         font_size=BODY_SIZE, font_color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

add_source(s10, 'Source: 舆情监测平台数据（2026-03-04 至 2026-03-11）')
add_page_number(s10, 10, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 11: 分角色行动建议 — 数据表格
# ───────────────────────────────────────────────────────────
s11 = prs.slides.add_slide(BL)
add_action_title(s11, '分角色行动建议：五大角色各有明确优先级，企业决策层与技术团队需紧急响应')

# 数据表格
headers11 = ['角色', '行动建议', '优先级']
col_widths11 = [Inches(2.5), Inches(7.0), Inches(2.233)]
hdr_y11 = Inches(1.5)
cx11 = LM
for hdr, cw in zip(headers11, col_widths11):
    add_text(s11, cx11, hdr_y11, cw, Inches(0.4), hdr,
             font_size=BODY_SIZE, font_color=MED_GRAY, bold=True)
    cx11 += cw
add_hline(s11, LM, Inches(2.0), CONTENT_W, BLACK, Pt(1.0))

roles = [
    ('企业决策层', '保持技术栈灵活性，V4发布首日测试复杂逻辑任务的成本效益比', '紧急'),
    ('技术团队', '重新评估RAG架构必要性，准备1M上下文接入方案', '紧急'),
    ('投资者', '关注真正完成V4适配的硬核标的，警惕纯情绪博弈', '重要'),
    ('内容创作者', '提前梳理DeepSeek技术演进路径，抢占评测内容先机', '重要'),
    ('算力供应商', '加速昇腾生态适配，抓住"华为+DeepSeek"标配化窗口', '紧急'),
]

row_h11 = Inches(0.8)
for ri, (role, action, priority) in enumerate(roles):
    ry = Inches(2.1) + row_h11 * ri
    # 角色
    add_text(s11, LM, ry, col_widths11[0], row_h11, role,
             font_size=BODY_SIZE, font_color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # 行动建议
    add_text(s11, LM + col_widths11[0], ry, col_widths11[1], row_h11, action,
             font_size=BODY_SIZE, font_color=DARK_GRAY,
             anchor=MSO_ANCHOR.MIDDLE)
    # 优先级
    p_color = NAVY if priority == '紧急' else MED_GRAY
    add_text(s11, LM + col_widths11[0] + col_widths11[1], ry, col_widths11[2], row_h11,
             priority, font_size=BODY_SIZE, font_color=p_color, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_hline(s11, LM, ry + row_h11, CONTENT_W, LINE_GRAY)

add_source(s11, 'Source: 战略分析团队建议，2026年3月11日')
add_page_number(s11, 11, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 12: Key Takeaway — 核心洞见页
# ───────────────────────────────────────────────────────────
s12 = prs.slides.add_slide(BL)
add_action_title(s12, 'V4 的意义已超越任何单一模型发布——它正在成为中国AI"算力主权化"的标志性事件')

# 左侧内容
add_text(s12, LM, Inches(1.4), Inches(7.5), Inches(0.4),
         '三维度战略价值', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
add_hline(s12, LM, Inches(1.9), Inches(7.5), LINE_GRAY)
add_text(s12, LM, Inches(2.1), Inches(7.5), Inches(4.0),
         ['技术维度：',
          '1M 上下文 + 原生多模态 + mHC/Engram 新架构，全面对标并超越国际顶尖水平。'
          'SWE-bench 83.7% 刷新纪录，推理成本可降低 90%。',
          '',
          '生态维度：',
          '战略性"去英伟达化"构建自主算力闭环，OpenClaw Agent框架增速超越Linux，'
          '"华为昇腾+DeepSeek"正在成为政企大模型标配。',
          '',
          '资本维度：',
          '"DeepSeek交易"成为全球半导体估值的宏观变量。V4发布将不仅改变AI竞争格局，'
          '更将深度影响全球科技资本流向。'],
         font_size=Pt(12), font_color=DARK_GRAY, line_spacing=Pt(4))

# 右侧 Key Takeaways
tk_x = Inches(9.0)
tk_w = Inches(3.533)
add_rect(s12, tk_x, Inches(1.4), tk_w, Inches(5.2), BG_GRAY)
add_text(s12, tk_x + Inches(0.2), Inches(1.6), tk_w - Inches(0.4), Inches(0.4),
         'Key Takeaways', font_size=BODY_SIZE, font_color=NAVY, bold=True)
add_hline(s12, tk_x + Inches(0.2), Inches(2.1), tk_w - Inches(0.4), LINE_GRAY)
add_text(s12, tk_x + Inches(0.2), Inches(2.3), tk_w - Inches(0.4), Inches(4.0),
         ['1. V4 重新定义中国AI"算力主权"路径',
          '',
          '2. 技术+生态+资本三维度联动，重塑2026年AI竞争格局底层逻辑',
          '',
          '3. 留给V4的"独占窗口期"有限，速度即战略',
          '',
          '4. 需警惕"利好出尽"风险，回归技术验证与API定价'],
         font_size=Pt(12), font_color=DARK_GRAY, line_spacing=Pt(4))

add_source(s12, 'Source: 综合分析，知识库监控报告（2026-03-04 至 2026-03-11）')
add_page_number(s12, 12, TOTAL_SLIDES)

# ───────────────────────────────────────────────────────────
#  Slide 13: Closing
# ───────────────────────────────────────────────────────────
s13 = prs.slides.add_slide(BL)
add_rect(s13, 0, 0, SW, Inches(0.05), NAVY)
add_text(s13, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
         '技术主权，算力自主，生态制胜', font_size=Pt(28), font_color=NAVY, bold=True,
         font_name='Georgia', alignment=PP_ALIGN.CENTER)
add_hline(s13, Inches(5.5), Inches(3.3), Inches(2.3), NAVY, Pt(1.5))
add_text(s13, Inches(1.5), Inches(3.8), Inches(10.3), Inches(2.0),
         'DeepSeek-V4 的意义已超越单一模型发布\n它正在成为中国AI产业"算力主权化"的标志性事件\n从技术、生态、资本三个维度重塑AI竞争格局',
         font_size=SUB_HEADER_SIZE, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
         line_spacing=Pt(14))
add_text(s13, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.4),
         'CONFIDENTIAL  |  战略分析报告  |  2026年3月', font_size=Pt(12),
         font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_hline(s13, Inches(1), Inches(6.8), Inches(3), NAVY, Pt(2))
add_page_number(s13, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════
#  Save & Cleanup
# ═══════════════════════════════════════════════════════════
outpath = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'output', 'DeepSeek-V4_发布动态周报.pptx')
os.makedirs(os.path.dirname(outpath), exist_ok=True)
prs.save(outpath)
full_cleanup(outpath)
print(f'✅ Presentation saved to: {outpath}')
