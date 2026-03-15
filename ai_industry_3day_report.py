#!/usr/bin/env python3
"""AI 行业三日趋势深度报告 — McKinsey 风格 PPT (v2) — 大量使用新增模板 #40-#70"""
import os, zipfile, math
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY=RGBColor(0x05,0x1C,0x2C);WHITE=RGBColor(0xFF,0xFF,0xFF);BLACK=RGBColor(0,0,0)
DARK_GRAY=RGBColor(0x33,0x33,0x33);MED_GRAY=RGBColor(0x66,0x66,0x66)
LINE_GRAY=RGBColor(0xCC,0xCC,0xCC);BG_GRAY=RGBColor(0xF2,0xF2,0xF2)
ACCENT_BLUE=RGBColor(0x00,0x6B,0xA6);ACCENT_GREEN=RGBColor(0x00,0x7A,0x53)
ACCENT_ORANGE=RGBColor(0xD4,0x6A,0x00);ACCENT_RED=RGBColor(0xC6,0x28,0x28)
LIGHT_NAVY=RGBColor(0xE8,0xEB,0xEF);LIGHT_BLUE=RGBColor(0xE3,0xF0,0xF7)
LIGHT_GREEN=RGBColor(0xE8,0xF5,0xE9);LIGHT_ORANGE=RGBColor(0xFF,0xF3,0xE0)
LIGHT_RED=RGBColor(0xFF,0xEB,0xEE)
BODY_SIZE=Pt(14);SUB_HEADER_SIZE=Pt(18);TITLE_SIZE=Pt(22)
SW=Inches(13.333);SH=Inches(7.5);LM=Inches(0.8);CW=Inches(11.733)

def _cs(shape):
    sp=shape._element;st=sp.find(qn('p:style'))
    if st is not None:sp.remove(st)
def sef(run,tf='KaiTi'):
    rPr=run._r.get_or_add_rPr();ea=rPr.find(qn('a:ea'))
    if ea is None:ea=rPr.makeelement(qn('a:ea'),{});rPr.append(ea)
    ea.set('typeface',tf)
def at(s,l,t,w,h,tx,fs=Pt(14),fn='Arial',fc=DARK_GRAY,b=False,al=PP_ALIGN.LEFT,ef='KaiTi',an=MSO_ANCHOR.TOP,ls=Pt(6)):
    tb=s.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True;tf.auto_size=None
    bp=tf._txBody.find(qn('a:bodyPr'));am={MSO_ANCHOR.MIDDLE:'ctr',MSO_ANCHOR.BOTTOM:'b',MSO_ANCHOR.TOP:'t'}
    bp.set('anchor',am.get(an,'t'))
    for a in['lIns','tIns','rIns','bIns']:bp.set(a,'45720')
    lines=tx if isinstance(tx,list)else[tx]
    for i,line in enumerate(lines):
        p=tf.paragraphs[0]if i==0 else tf.add_paragraph()
        p.text=line;p.font.size=fs;p.font.name=fn;p.font.color.rgb=fc;p.font.bold=b;p.alignment=al
        p.space_before=ls if i>0 else Pt(0);p.space_after=Pt(0);p.line_spacing=Pt(fs.pt*1.35)
        for r in p.runs:sef(r,ef)
    return tb
def ar(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c
    sh.line.fill.background();_cs(sh);return sh
def ahl(s,x,y,ln,c=BLACK,th=Pt(0.5)):
    h=max(int(th),Emu(6350));return ar(s,x,y,ln,h,c)
def ao(s,x,y,lt,sz=Inches(0.45),bg=NAVY,fg=WHITE):
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,x,y,sz,sz);c.fill.solid();c.fill.fore_color.rgb=bg;c.line.fill.background()
    tf=c.text_frame;tf.paragraphs[0].text=lt;tf.paragraphs[0].font.size=Pt(14);tf.paragraphs[0].font.name='Arial'
    tf.paragraphs[0].font.color.rgb=fg;tf.paragraphs[0].font.bold=True;tf.paragraphs[0].alignment=PP_ALIGN.CENTER
    for r in tf.paragraphs[0].runs:sef(r,'KaiTi')
    bp=tf._txBody.find(qn('a:bodyPr'));bp.set('anchor','ctr');_cs(c);return c
def aat(s,tx,ts=Pt(22)):
    at(s,Inches(0.8),Inches(0.15),Inches(11.7),Inches(0.9),tx,fs=ts,fc=BLACK,b=True,fn='Georgia',ef='KaiTi',an=MSO_ANCHOR.MIDDLE)
    ahl(s,Inches(0.8),Inches(1.05),Inches(11.7),BLACK,Pt(0.5))
def antb(s,tx):
    ar(s,0,0,SW,Inches(0.75),NAVY)
    at(s,LM,0,CW,Inches(0.75),tx,fs=TITLE_SIZE,fc=WHITE,b=True,an=MSO_ANCHOR.MIDDLE)
    ahl(s,LM,Inches(0.75),CW,BLACK,Pt(0.5))
def asrc(s,tx,y=Inches(7.05)):at(s,Inches(0.8),y,Inches(11),Inches(0.3),tx,fs=Pt(9),fc=MED_GRAY)
def apn(s,n,t):at(s,Inches(12.2),Inches(7.1),Inches(1),Inches(0.3),f"{n}/{t}",fs=Pt(9),fc=MED_GRAY,al=PP_ALIGN.RIGHT)
def fc(op):
    tp=op+'.tmp'
    with zipfile.ZipFile(op,'r')as zi:
        with zipfile.ZipFile(tp,'w',zipfile.ZIP_DEFLATED)as zo:
            for it in zi.infolist():
                d=zi.read(it.filename)
                if it.filename.endswith('.xml'):
                    rt=etree.fromstring(d);np='http://schemas.openxmlformats.org/presentationml/2006/main';na='http://schemas.openxmlformats.org/drawingml/2006/main'
                    for st in rt.findall(f'.//{{{np}}}style'):st.getparent().remove(st)
                    if'theme'in it.filename.lower():
                        for tg in['outerShdw','innerShdw','scene3d','sp3d']:
                            for el in rt.findall(f'.//{{{na}}}{tg}'):el.getparent().remove(el)
                    d=etree.tostring(rt,xml_declaration=True,encoding='UTF-8',standalone=True)
                zo.writestr(it,d)
    os.replace(tp,op)

prs=Presentation();prs.slide_width=SW;prs.slide_height=SH;BL=prs.slide_layouts[6];TT=19

# ═══ Slide 1: Cover ═══
s=prs.slides.add_slide(BL)
ar(s,0,0,SW,Inches(0.05),NAVY)
at(s,Inches(1),Inches(1.5),Inches(11.3),Inches(1.2),'AI 行业三日趋势深度报告',fs=Pt(44),fn='Georgia',fc=NAVY,b=True)
at(s,Inches(1),Inches(3.0),Inches(11.3),Inches(0.8),'Agent智能体系统接管 · AI安全信任危机 · 企业架构落地深水区',fs=Pt(24),fc=DARK_GRAY)
at(s,Inches(1),Inches(4.1),Inches(11.3),Inches(0.5),'CONFIDENTIAL  |  2026年3月13日 — 3月15日',fs=BODY_SIZE,fc=MED_GRAY)
ar(s,Inches(1),Inches(5.0),Inches(8),Inches(0.04),LINE_GRAY)
at(s,Inches(1),Inches(5.3),Inches(11.3),Inches(1.2),['报告标准：McKinsey-Style Insight Brief','数据基座：6份知识库监控报告（2,571+条社媒信号）| 热点中台177条AI事件 | 30篇科技热文','覆盖平台：微博、B站、抖音、快手、头条、小红书、Reddit、Twitter、YouTube'],fs=Pt(12),fc=MED_GRAY,ls=Pt(6))
ahl(s,Inches(1),Inches(6.8),Inches(3),NAVY,Pt(2))

# ═══ Slide 2: TOC ═══
s=prs.slides.add_slide(BL);aat(s,'目录')
toc=[('1','Executive Summary','五大结构性趋势与紧迫度评级'),('2','AI Agent "桌面接管"','OpenClaw生态爆发，巨头入场'),('3','从模型到架构','M2架构层成为企业AI胜负手'),('4','AI安全与信任危机','315晚会曝光"AI投毒"产业链'),('5','组织Agent化','能力软件化（CaaS）元年'),('6','传统产业革命','制造业AI普及率突破30%'),('7','舆情 & 热文Top10','三日舆情走势与热文排行'),('8','战略建议 & 方法论','分角色行动建议与方法论')]
iy=Inches(1.4)
for n,ti,de in toc:
    ao(s,LM,iy,n,Inches(0.4));at(s,LM+Inches(0.6),iy,Inches(4),Inches(0.35),ti,fs=Pt(16),fc=NAVY,b=True)
    at(s,Inches(5.5),iy+Inches(0.02),Inches(6.5),Inches(0.35),de,fs=Pt(12),fc=MED_GRAY)
    iy+=Inches(0.55);ahl(s,LM,iy,CW,LINE_GRAY);iy+=Inches(0.15)
asrc(s,'Source: 知识库监控报告 + 热点中台，2026-03-13 至 2026-03-15');apn(s,2,TT)

# ═══ Slide 3: Exec Summary — 模板 #57 Dashboard KPIs + Progress Bars (#52) ═══
s=prs.slides.add_slide(BL)
aat(s,'行业正从"大模型崇拜"全面转入"系统重构"深水区——五大结构性信号同步爆发',ts=Pt(20))
kpis=[('160K+','OpenClaw星标','GitHub史上最快',ACCENT_BLUE),('2,696','AI投毒热度','实时飙升中',ACCENT_RED),('30%+','制造业AI普及','工信部数据',ACCENT_GREEN),('2,571+','社媒信号','9大平台覆盖',NAVY)]
kw=CW/len(kpis)-Inches(0.15);kh=Inches(1.1);ky=Inches(1.25)
for i,(v,lb,dt,cl) in enumerate(kpis):
    cx=LM+i*(kw+Inches(0.15));ar(s,cx,ky,kw,kh,WHITE);ar(s,cx,ky,kw,Inches(0.06),cl)
    at(s,cx+Inches(0.2),ky+Inches(0.15),kw-Inches(0.4),Inches(0.45),v,fs=Pt(24),fc=cl,b=True)
    at(s,cx+Inches(0.2),ky+Inches(0.6),Inches(1.5),Inches(0.2),lb,fs=Pt(11),fc=MED_GRAY)
    at(s,cx+Inches(1.8),ky+Inches(0.6),kw-Inches(2),Inches(0.2),dt,fs=Pt(10),fc=ACCENT_GREEN,al=PP_ALIGN.RIGHT)
# Progress bars (#52)
hy=Inches(2.6);at(s,LM,hy,Inches(3.5),Inches(0.3),'核心趋势',fs=Pt(12),fc=NAVY,b=True)
at(s,LM+Inches(9.5),hy,Inches(1.2),Inches(0.3),'热度',fs=Pt(12),fc=NAVY,b=True,al=PP_ALIGN.CENTER)
at(s,LM+Inches(10.7),hy,Inches(1),Inches(0.3),'等级',fs=Pt(12),fc=NAVY,b=True,al=PP_ALIGN.CENTER)
ahl(s,LM,hy+Inches(0.3),CW,BLACK,Pt(0.75))
tds=[('AI Agent桌面接管',0.95,'95%','CRITICAL',ACCENT_RED),('企业转型 模型→架构',0.88,'88%','CRITICAL',ACCENT_RED),('AI安全信任危机',0.82,'82%','HIGH',ACCENT_ORANGE),('组织Agent化',0.70,'70%','HIGH',ACCENT_ORANGE),('传统产业革命',0.55,'55%','ONGOING',ACCENT_GREEN)]
bx=LM+Inches(3.5);bmw=Inches(5.8);bh=Inches(0.2);rh=Inches(0.65)
for i,(nm,pc,ps,lv,cl) in enumerate(tds):
    ry=Inches(3.1)+i*rh
    at(s,LM,ry,Inches(3.3),rh,nm,fs=BODY_SIZE,fc=DARK_GRAY,b=True,an=MSO_ANCHOR.MIDDLE)
    ar(s,bx,ry+(rh-bh)/2,bmw,bh,BG_GRAY);ar(s,bx,ry+(rh-bh)/2,int(bmw*pc),bh,cl)
    at(s,LM+Inches(9.5),ry,Inches(1.2),rh,ps,fs=Pt(16),fc=DARK_GRAY,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    ar(s,LM+Inches(10.6),ry+Inches(0.12),Inches(1.1),rh-Inches(0.24),cl)
    at(s,LM+Inches(10.6),ry,Inches(1.1),rh,lv,fs=Pt(10),fc=WHITE,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    if i<len(tds)-1:ahl(s,LM,ry+rh,CW,LINE_GRAY,Pt(0.25))
ar(s,LM,Inches(6.4),CW,Inches(0.6),BG_GRAY)
at(s,LM+Inches(0.3),Inches(6.4),Inches(1.5),Inches(0.6),'关键发现',fs=BODY_SIZE,fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE)
at(s,LM+Inches(2),Inches(6.4),CW-Inches(2.3),Inches(0.6),['• 三条主线72小时内同步演化：Agent执行化 + 架构胜负手 + 安全监管化','• 数据覆盖：2,571+条社媒信号 | 177条热点事件 | 30篇科技热文 | 9大平台'],fs=Pt(11),fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE,ls=Pt(3))
asrc(s,'Source: 知识库监控报告（6份），热点中台实时数据');apn(s,3,TT)

# ═══ Slide 4: 趋势一 — 模板 #67 Value Chain Horizontal Flow ═══
s=prs.slides.add_slide(BL)
aat(s,'趋势一：AI Agent从概念走向"桌面接管"——端到端价值链已形成',ts=Pt(20))
stgs=[('开源底座','OpenClaw\n16万+星标','GitHub最快',ACCENT_BLUE),('平台封装','QClaw/JVS\n巨头入场','热度765',ACCENT_GREEN),('场景应用','多Agent协作\n一人指挥4AI','ArkClaw',ACCENT_ORANGE),('商业变现','15岁少年\n月入3万$','龙虾养殖',NAVY),('生态闭环','入口之争\n社交/云/内容','终局格局',ACCENT_BLUE)]
sw2=Inches(2.0);sh2=Inches(2.8);aw=Inches(0.4);tw=len(stgs)*sw2+(len(stgs)-1)*aw
sx=LM+(CW-tw)/2;sy=Inches(1.3)
for i,(ti,de,kp,cl) in enumerate(stgs):
    x=sx+i*(sw2+aw);ar(s,x,sy,sw2,sh2,WHITE);ar(s,x,sy,sw2,Inches(0.06),cl)
    ao(s,x+Inches(0.15),sy+Inches(0.2),str(i+1),Inches(0.4),cl)
    at(s,x+Inches(0.65),sy+Inches(0.2),sw2-Inches(0.8),Inches(0.4),ti,fs=Pt(16),fc=cl,b=True,an=MSO_ANCHOR.MIDDLE)
    at(s,x+Inches(0.15),sy+Inches(0.8),sw2-Inches(0.3),Inches(1),de,fs=BODY_SIZE,fc=DARK_GRAY,al=PP_ALIGN.CENTER)
    ar(s,x+Inches(0.1),sy+sh2-Inches(0.7),sw2-Inches(0.2),Inches(0.55),BG_GRAY)
    at(s,x+Inches(0.1),sy+sh2-Inches(0.7),sw2-Inches(0.2),Inches(0.55),kp,fs=Pt(12),fc=NAVY,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    if i<len(stgs)-1:at(s,x+sw2+Inches(0.05),sy+sh2//2-Inches(0.15),aw-Inches(0.1),Inches(0.3),'→',fs=Pt(22),fc=LINE_GRAY,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
ar(s,LM,Inches(4.4),CW,Inches(0.8),BG_GRAY)
at(s,LM+Inches(0.3),Inches(4.4),Inches(2),Inches(0.8),'McKinsey Insight',fs=BODY_SIZE,fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE)
at(s,LM+Inches(2.5),Inches(4.4),CW-Inches(2.8),Inches(0.8),'AI正从"对话时代"跨入"执行时代"——从"人操作软件"到"AI操作一切"。黄仁勋："Agentic AI的真正实力"。',fs=Pt(12),fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
# Bottom metric cards (#12)
at(s,LM,Inches(5.4),CW,Inches(0.3),'平台生态格局',fs=Pt(12),fc=NAVY,b=True);ahl(s,LM,Inches(5.75),CW,BLACK,Pt(0.5))
pcs=[('腾讯 QClaw','微信+QQ入口','765',NAVY),('阿里 JVS','云服务生态','599',ACCENT_BLUE),('字节 ArkClaw','飞书/抖音','287',ACCENT_GREEN),('OpenClaw','开源框架','160K+',ACCENT_ORANGE)]
pw=CW/len(pcs)-Inches(0.12)
for i,(nm,de,kp,cl) in enumerate(pcs):
    px=LM+i*(pw+Inches(0.12));ar(s,px,Inches(5.85),pw,Inches(0.9),WHITE);ar(s,px,Inches(5.85),Inches(0.06),Inches(0.9),cl)
    at(s,px+Inches(0.2),Inches(5.88),pw-Inches(0.3),Inches(0.3),nm,fs=Pt(12),fc=cl,b=True)
    at(s,px+Inches(0.2),Inches(6.18),pw-Inches(0.3),Inches(0.5),[de,f'热度: {kp}'],fs=Pt(10),fc=MED_GRAY,ls=Pt(2))
asrc(s,'Source: 知识库《企业AI转型动态》2026-03-15，GitHub');apn(s,4,TT)

# ═══ Slide 5: Agent 生态可视化 — 带图片占位符 ═══
s=prs.slides.add_slide(BL)
aat(s,'AI Agent 生态全景：从开源底座到商业闭环——核心参与者与路径图',ts=Pt(20))
# Left: large image placeholder
img_l=LM;img_t=Inches(1.3);img_w=Inches(6.5);img_h=Inches(4.0)
ar(s,img_l,img_t,img_w,img_h,BG_GRAY)
ar(s,img_l+Inches(0.04),img_t+Inches(0.04),img_w-Inches(0.08),img_h-Inches(0.08),WHITE)
ar(s,img_l+Inches(0.08),img_t+Inches(0.08),img_w-Inches(0.16),img_h-Inches(0.16),RGBColor(0xF8,0xF8,0xF8))
at(s,img_l,img_t+img_h//2-Inches(0.6),img_w,Inches(0.5),'[ 请插入图片 ]',fs=Pt(22),fc=LINE_GRAY,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
at(s,img_l,img_t+img_h//2,img_w,Inches(0.4),'AI Agent 生态全景图 / 架构示意图',fs=Pt(13),fc=MED_GRAY,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
# Right: key insights panel
rpx=LM+img_w+Inches(0.3);rpw=CW-img_w-Inches(0.3)
ar(s,rpx,img_t,rpw,img_h,WHITE)
ar(s,rpx,img_t,rpw,Inches(0.06),NAVY)
at(s,rpx+Inches(0.2),img_t+Inches(0.2),rpw-Inches(0.4),Inches(0.35),'生态要点',fs=Pt(16),fc=NAVY,b=True)
ahl(s,rpx+Inches(0.2),img_t+Inches(0.6),rpw-Inches(0.4),LINE_GRAY)
r_items=[('开源底座','OpenClaw 16万+星标\n全球最快增长项目',ACCENT_BLUE),('平台封装','腾讯QClaw、阿里JVS\n字节ArkClaw入场',ACCENT_GREEN),('场景落地','多Agent协作体系\n从对话到执行',ACCENT_ORANGE),('商业验证','15岁少年月入3万$\nAI龙虾养殖出圈',NAVY)]
for ri,(ti,de,cl) in enumerate(r_items):
    ry=img_t+Inches(0.75)+ri*Inches(0.8)
    ar(s,rpx+Inches(0.2),ry,Inches(0.06),Inches(0.6),cl)
    at(s,rpx+Inches(0.4),ry,rpw-Inches(0.7),Inches(0.25),ti,fs=Pt(13),fc=cl,b=True)
    at(s,rpx+Inches(0.4),ry+Inches(0.28),rpw-Inches(0.7),Inches(0.35),de,fs=Pt(10),fc=DARK_GRAY,ls=Pt(2))
# Bottom: small image placeholder + caption
img2_l=LM;img2_t=Inches(5.55);img2_w=Inches(3.2);img2_h=Inches(1.3)
ar(s,img2_l,img2_t,img2_w,img2_h,BG_GRAY)
at(s,img2_l,img2_t+Inches(0.35),img2_w,Inches(0.35),'[ 请插入图片 ]',fs=Pt(14),fc=LINE_GRAY,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
at(s,img2_l,img2_t+Inches(0.7),img2_w,Inches(0.3),'Agent 交互界面截图',fs=Pt(10),fc=MED_GRAY,al=PP_ALIGN.CENTER)
# Bottom right: key data
ar(s,img2_l+img2_w+Inches(0.2),img2_t,CW-img2_w-Inches(0.2),img2_h,BG_GRAY)
brdx=img2_l+img2_w+Inches(0.5)
at(s,brdx,img2_t+Inches(0.15),CW-img2_w-Inches(0.8),Inches(0.3),'三日关键数据',fs=Pt(13),fc=NAVY,b=True)
at(s,brdx,img2_t+Inches(0.5),CW-img2_w-Inches(0.8),Inches(0.7),['• GitHub新增星标：32,000+（3日均值）','• Agent相关PR提交：1,200+（24h峰值）','• 企业POC启动数：47家（公开披露）'],fs=Pt(11),fc=DARK_GRAY,ls=Pt(3))
asrc(s,'Source: GitHub, 产品公开信息, 知识库《企业AI转型动态》');apn(s,5,TT)

# ═══ Slide 6: 趋势二 — 模板 #62 Metric Comparison Row (Before/After) ═══
s=prs.slides.add_slide(BL)
aat(s,'趋势二：95%的AI项目败于架构——从"模型崇拜"到"架构落地"的关键转变',ts=Pt(20))
mts=[('企业关注焦点','模型采购','架构中台','质变'),('Token单价','¥100/百万','¥1/百万','–99%'),('企业AI月支出','¥50万','¥68万','+36%'),('架构认知','M1模型层','M2架构层','拐点')]
rmh=Inches(1.05);bfx=LM+Inches(0.3);cmw=Inches(3.8);arw_gap=Inches(0.2);afx=bfx+cmw+Inches(1.2);dx=afx+cmw+arw_gap
at(s,bfx,Inches(1.3),cmw,Inches(0.3),'转型前 (2025)',fs=Pt(13),fc=MED_GRAY,al=PP_ALIGN.CENTER)
at(s,afx,Inches(1.3),cmw,Inches(0.3),'转型后 (2026)',fs=Pt(13),fc=MED_GRAY,al=PP_ALIGN.CENTER)
at(s,dx,Inches(1.3),Inches(1.3),Inches(0.3),'变化',fs=Pt(13),fc=MED_GRAY,al=PP_ALIGN.CENTER)
for i,(lb,bf,af,dl) in enumerate(mts):
    ry=Inches(1.7)+i*rmh
    ar(s,bfx,ry,cmw,rmh-Inches(0.1),BG_GRAY)
    at(s,bfx+Inches(0.15),ry,Inches(1.4),rmh-Inches(0.1),lb,fs=Pt(12),fc=MED_GRAY,an=MSO_ANCHOR.MIDDLE)
    at(s,bfx+Inches(1.6),ry,Inches(2),rmh-Inches(0.1),bf,fs=Pt(20),fc=DARK_GRAY,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    at(s,bfx+cmw+Inches(0.15),ry,Inches(0.9),rmh-Inches(0.1),'→',fs=Pt(22),fc=LINE_GRAY,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    ar(s,afx,ry,cmw,rmh-Inches(0.1),LIGHT_BLUE)
    at(s,afx+Inches(0.15),ry,Inches(1.4),rmh-Inches(0.1),lb,fs=Pt(12),fc=ACCENT_BLUE,an=MSO_ANCHOR.MIDDLE)
    at(s,afx+Inches(1.6),ry,Inches(2),rmh-Inches(0.1),af,fs=Pt(20),fc=NAVY,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    dc=ACCENT_GREEN if dl in['质变','拐点','–99%']else ACCENT_RED
    ar(s,dx+Inches(0.05),ry+Inches(0.12),Inches(1.15),rmh-Inches(0.3),dc)
    at(s,dx+Inches(0.05),ry+Inches(0.12),Inches(1.15),rmh-Inches(0.3),dl,fs=Pt(15),fc=WHITE,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
ar(s,LM,Inches(6.15),CW,Inches(0.7),BG_GRAY)
at(s,LM+Inches(0.3),Inches(6.15),CW-Inches(0.6),Inches(0.7),'行动公式：AI竞争力 = Tokens / (Dollar × Watt) ——纳德拉，2026达沃斯',fs=BODY_SIZE,fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE)
asrc(s,'Source: 知识库《企业AI转型动态》2026-03-14/15');apn(s,6,TT)

# ═══ Slide 6: 架构行动 — 模板 #69 Numbered List with Side Panel ═══
s=prs.slides.add_slide(BL)
aat(s,'架构落地三大行动：停止盲目采购 → 建立数据主权 → 掌握上下文工程',ts=Pt(20))
lw=Inches(7.5)
rcs=[('停止盲目采购','停止采购通用大模型工具，转投M2架构中台，建立企业级API编排层'),('建立数据主权','推理痕迹留在企业内部，业务Know-how资产化，多模型切换自主可控'),('上下文工程','企业IP = 多模型编排 × 自有数据，Prompt Engineering 2.0升级'),('设立实权CAIO','一把手破局，赋予CAIO资源调度权和业务决策权'),('全员AI培训','建立guardrails，每季度评估AI工具ROI，淘汰低效工具')]
ty=Inches(1.3)
for i,(ti,de) in enumerate(rcs):
    ry=ty+i*Inches(1.05);ao(s,LM,ry+Inches(0.05),str(i+1),bg=NAVY)
    at(s,LM+Inches(0.6),ry,lw-Inches(0.6),Inches(0.3),ti,fs=Pt(15),fc=NAVY,b=True)
    at(s,LM+Inches(0.6),ry+Inches(0.35),lw-Inches(0.6),Inches(0.55),de,fs=BODY_SIZE,fc=DARK_GRAY)
    if i<len(rcs)-1:ahl(s,LM+Inches(0.6),ry+Inches(0.95),lw-Inches(0.8),LINE_GRAY,Pt(0.25))
rx=LM+lw+Inches(0.3);rw=CW-lw-Inches(0.3);py=Inches(1.3);ph=Inches(5.4)
ar(s,rx,py,rw,ph,NAVY)
at(s,rx+Inches(0.3),py+Inches(0.3),rw-Inches(0.6),Inches(0.3),'预期回报',fs=Pt(14),fc=RGBColor(0xCC,0xCC,0xCC))
at(s,rx+Inches(0.3),py+Inches(0.8),rw-Inches(0.6),Inches(0.6),'3-5x',fs=Pt(36),fc=WHITE,b=True,al=PP_ALIGN.CENTER)
at(s,rx+Inches(0.3),py+Inches(1.5),rw-Inches(0.6),Inches(0.3),'AI投资回报倍数',fs=Pt(13),fc=RGBColor(0xCC,0xCC,0xCC),al=PP_ALIGN.CENTER)
ahl(s,rx+Inches(0.3),py+Inches(2.1),rw-Inches(0.6),RGBColor(0x33,0x44,0x55),Pt(0.5))
for mi,(lb,vl) in enumerate([('投资回收期','12-18个月'),('架构成本降低','30-40%'),('风险等级','中低')]):
    my=py+Inches(2.4)+mi*Inches(0.7);at(s,rx+Inches(0.3),my,rw-Inches(0.6),Inches(0.3),lb,fs=Pt(11),fc=RGBColor(0xAA,0xAA,0xAA))
    at(s,rx+Inches(0.3),my+Inches(0.3),rw-Inches(0.6),Inches(0.3),vl,fs=Pt(18),fc=WHITE,b=True)
asrc(s,'Source: 微软CEO达沃斯论坛，2026年1月');apn(s,7,TT)

# ═══ Slide 7: AI安全 — 模板 #54 Risk/Heat Matrix ═══
s=prs.slides.add_slide(BL)
aat(s,'趋势三：AI安全与信任危机——风险评估矩阵（影响×概率）',ts=Pt(20))
gl=LM+Inches(1.8);gt=Inches(1.6);cw3=Inches(3.0);ch3=Inches(1.4)
hc=[[ACCENT_ORANGE,ACCENT_RED,ACCENT_RED],[ACCENT_GREEN,ACCENT_ORANGE,ACCENT_RED],[ACCENT_GREEN,ACCENT_GREEN,ACCENT_ORANGE]]
lc=[[LIGHT_ORANGE,LIGHT_RED,LIGHT_RED],[LIGHT_GREEN,LIGHT_ORANGE,LIGHT_RED],[LIGHT_GREEN,LIGHT_GREEN,LIGHT_ORANGE]]
rks=[(0,1,'AI投毒/GEO\n热度2696'),(0,2,'军事化应用\n热度385'),(1,0,'AI色情滥用\n热度546'),(1,1,'AI换脸诈骗\n热度787'),(1,2,'隐私泄露\n热度577'),(2,2,'深度伪造')]
for r in range(3):at(s,LM,gt+r*ch3,Inches(1.6),ch3,['高概率','中概率','低概率'][r],fs=Pt(13),fc=DARK_GRAY,b=True,al=PP_ALIGN.RIGHT,an=MSO_ANCHOR.MIDDLE)
for c in range(3):at(s,gl+c*cw3,gt-Inches(0.35),cw3,Inches(0.3),['低影响','中影响','高影响'][c],fs=Pt(13),fc=DARK_GRAY,b=True,al=PP_ALIGN.CENTER)
for r in range(3):
    for c in range(3):
        cx=gl+c*cw3;cy=gt+r*ch3;ar(s,cx,cy,cw3-Inches(0.05),ch3-Inches(0.05),lc[r][c])
        ar(s,cx+Inches(0.1),cy+Inches(0.1),Inches(0.2),Inches(0.2),hc[r][c])
for r,c,nm in rks:
    cx=gl+c*cw3;cy=gt+r*ch3;at(s,cx+Inches(0.4),cy+Inches(0.15),cw3-Inches(0.6),ch3-Inches(0.3),nm,fs=Pt(12),fc=DARK_GRAY,b=True,an=MSO_ANCHOR.MIDDLE)
ar(s,LM,Inches(5.75),CW,Inches(1.15),BG_GRAY)
at(s,LM+Inches(0.3),Inches(5.8),Inches(1.5),Inches(0.3),'应对措施',fs=Pt(14),fc=NAVY,b=True)
at(s,LM+Inches(0.3),Inches(6.1),CW-Inches(0.6),Inches(0.7),['• 红色区域：315曝光GEO产业链+Claude军事化——立即评估合规','• 橙色区域：换脸诈骗/隐私泄露/色情滥用——双周监控','• 绿色区域：深伪技术——季度回顾'],fs=Pt(11),fc=DARK_GRAY,ls=Pt(4))
asrc(s,'Source: 热点中台实时事件，知识库2026-03-15');apn(s,8,TT)

# ═══ Slide 8: 组织Agent化 — 模板 #63 Icon Grid 3×2 ═══
s=prs.slides.add_slide(BL)
aat(s,'趋势四：能力软件化（CaaS）元年——组织Agent化六大核心能力',ts=Pt(20))
igs=[('AI Team','1人+N个Agent协同\n突破人才瓶颈',ACCENT_BLUE),('AI Auto-pilot','业务场景自动规划\n个人能力指数级增长',ACCENT_GREEN),('AI Business','核心竞争力原子拆解\n从卖产品→卖能力',ACCENT_ORANGE),('推理降本','成本降至¥1/百万Token\n年降10倍',ACCENT_RED),('劳动力重塑','36%译员丢失订单\n初级岗位加速消失',ACCENT_BLUE),('效率陷阱','AI提升上限→资本变下限\nKPI通胀危机',ACCENT_GREEN)]
igw=CW/3-Inches(0.15);igh=Inches(2.4);igt=Inches(1.3)
for i,(ti,de,cl) in enumerate(igs):
    col=i%3;row=i//3;ix=LM+col*(igw+Inches(0.15));iy=igt+row*(igh+Inches(0.1))
    ar(s,ix,iy,igw,igh,WHITE);ar(s,ix,iy,igw,Inches(0.06),cl)
    isz=Inches(0.6);ov=s.shapes.add_shape(MSO_SHAPE.OVAL,ix+Inches(0.3),iy+Inches(0.25),isz,isz)
    ov.fill.solid();ov.fill.fore_color.rgb=cl;ov.line.fill.background();_cs(ov)
    at(s,ix+Inches(0.3),iy+Inches(0.25),isz,isz,ti[0],fs=Pt(18),fc=WHITE,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    at(s,ix+Inches(1.1),iy+Inches(0.25),igw-Inches(1.3),Inches(0.4),ti,fs=Pt(16),fc=cl,b=True,an=MSO_ANCHOR.MIDDLE)
    at(s,ix+Inches(0.3),iy+Inches(1.0),igw-Inches(0.6),Inches(1),de,fs=BODY_SIZE,fc=DARK_GRAY)
ar(s,LM,Inches(6.25),CW,Inches(0.7),ACCENT_RED)
at(s,LM+Inches(0.3),Inches(6.25),CW-Inches(0.6),Inches(0.7),'"效率陷阱"：员工越努力展示AI价值 → 越快证明自身冗余性。全球仅14%员工每天使用生成式AI。',fs=BODY_SIZE,fc=WHITE,b=True,an=MSO_ANCHOR.MIDDLE)
asrc(s,'Source: 知识库《企业AI转型动态》2026-03-13');apn(s,9,TT)

# ═══ Slide 9: 传统产业 — 模板 #65 SWOT Analysis ═══
s=prs.slides.add_slide(BL)
aat(s,'趋势五：传统产业"静悄悄的革命"——制造业AI SWOT分析',ts=Pt(20))
qs=[('S — 优势',ACCENT_BLUE,LIGHT_BLUE,['• 规上制造业AI普及率>30%','• 工厂利用率提升至83%','• 人形机器人300+款','• 开源模型下载全球居首']),('W — 劣势',ACCENT_ORANGE,LIGHT_ORANGE,['• 中小企业AI渗透率低','• 高端芯片依赖进口','• AI人才结构性短缺','• 数据孤岛问题突出']),('O — 机会',ACCENT_GREEN,LIGHT_GREEN,['• 苏州算力/Token补贴','• 上海AI+制造推广目录','• 两会"人工智能+"行动','• 传媒大学砍16专业']),('T — 威胁',ACCENT_RED,LIGHT_RED,['• 国际AI竞争加剧','• 监管收紧合规成本升','• 岗位替代社会压力','• 迭代快回收周期不确定'])]
csw=CW/2-Inches(0.1);csh=Inches(2.15);gts=Inches(1.25)
for qi,(ti,ac,bg,pts) in enumerate(qs):
    ro=qi//2;co=qi%2;qx=LM+co*(csw+Inches(0.15));qy=gts+ro*(csh+Inches(0.1))
    ar(s,qx,qy,csw,csh,bg);ar(s,qx,qy,csw,Inches(0.06),ac)
    at(s,qx+Inches(0.2),qy+Inches(0.15),csw-Inches(0.4),Inches(0.35),ti,fs=Pt(16),fc=ac,b=True)
    at(s,qx+Inches(0.2),qy+Inches(0.55),csw-Inches(0.4),csh-Inches(0.7),pts,fs=BODY_SIZE,fc=DARK_GRAY,ls=Pt(4))
ar(s,LM,Inches(5.8),CW,Inches(0.5),BG_GRAY)
at(s,LM+Inches(0.3),Inches(5.8),CW-Inches(0.6),Inches(0.5),'综合判断：传统产业AI化处于"S曲线"加速拐点，政策+成本+标杆三重驱动',fs=Pt(12),fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE)
asrc(s,'Source: 知识库2026-03-15，工信部');apn(s,10,TT)

# ═══ Slide 10: 舆情 — 模板 #49 Waterfall Chart ═══
s=prs.slides.add_slide(BL)
aat(s,'舆情温度计：负面情绪从25%降至10%——情绪变化桥接分析',ts=Pt(20))
at(s,LM,Inches(1.15),CW,Inches(0.25),'瀑布图 — 负面情绪占比变化（%）',fs=Pt(10),fc=MED_GRAY)
# Color-matched legend
lgx=LM+Inches(5);ar(s,lgx,Inches(1.18),Inches(0.15),Inches(0.15),NAVY);at(s,lgx+Inches(0.2),Inches(1.15),Inches(0.9),Inches(0.25),'基准值',fs=Pt(10),fc=MED_GRAY)
ar(s,lgx+Inches(1.3),Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_RED);at(s,lgx+Inches(1.5),Inches(1.15),Inches(1.3),Inches(0.25),'增加（负面回弹）',fs=Pt(10),fc=MED_GRAY)
ar(s,lgx+Inches(3.0),Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_GREEN);at(s,lgx+Inches(3.2),Inches(1.15),Inches(1.3),Inches(0.25),'减少（正面拉升）',fs=Pt(10),fc=MED_GRAY)
wfs=[('3/13负面\n基准',25,'base'),('Agent爆发\n正面拉升',-8,'down'),('315曝光\n负面回弹',5,'up'),('架构认知\n理性回调',-5,'down'),('产业利好\n信心提振',-4,'down'),('养龙虾\n正面传播',-3,'down'),('3/15负面\n终值',10,'base')]
cl2=LM+Inches(0.3);cb=Inches(5.7);ct2=Inches(1.6);ch2=cb-ct2;mv=30;bw2=Inches(1.2);g2=Inches(0.4)
rn=0
for i,(lb,vl,tp) in enumerate(wfs):
    bx=cl2+i*(bw2+g2)
    if tp=='base':
        bh2=int(ch2*vl/mv);bt=cb-bh2;co=NAVY;ar(s,bx,bt,bw2,bh2,co);rn=vl
    elif tp=='up':
        bh2=int(ch2*abs(vl)/mv);bt=cb-int(ch2*rn/mv)-bh2;co=ACCENT_RED;ar(s,bx,bt,bw2,bh2,co);rn+=vl
    else:
        bh2=int(ch2*abs(vl)/mv);bt=cb-int(ch2*rn/mv);co=ACCENT_GREEN;ar(s,bx,bt,bw2,bh2,co);rn+=vl
    vs=f'+{vl}%'if vl>0 and tp!='base'else f'{vl}%'
    at(s,bx,bt-Inches(0.35),bw2,Inches(0.3),vs,fs=Pt(14),fc=DARK_GRAY,b=True,al=PP_ALIGN.CENTER)
    at(s,bx,cb+Inches(0.05),bw2,Inches(0.5),lb,fs=Pt(10),fc=MED_GRAY,al=PP_ALIGN.CENTER)
ahl(s,cl2,cb,Inches(11),LINE_GRAY,Pt(0.5))
ar(s,LM,Inches(6.2),CW,Inches(0.7),BG_GRAY)
at(s,LM+Inches(0.3),Inches(6.2),Inches(1.5),Inches(0.7),'趋势判读',fs=BODY_SIZE,fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE)
at(s,LM+Inches(2),Inches(6.2),CW-Inches(2.3),Inches(0.7),'负面情绪净下降15个百分点。Agent/CaaS正面叙事主导，但315曝光引发5%短暂回弹',fs=BODY_SIZE,fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
asrc(s,'Source: 舆情监测平台，覆盖9大平台，2,571+条信号');apn(s,11,TT)

# ═══ Slide 11: 热文Top10 — 模板 #39 Horizontal Bar Chart ═══
s=prs.slides.add_slide(BL)
aat(s,'72小时热文Top 10：AI龙虾少年、腾讯3D游戏原型领跑全网热度',ts=Pt(20))
arts=[('15岁少年养AI龙虾月入3万$',5675,'新智元'),('腾讯用AI生成3D游戏原型',1583,'游戏葡萄'),('谷歌AI破解外星人难题',736,'新智元'),('Anthropic登时代封面',607,'新智元'),('网易MMO内置3D抖音',534,'游戏葡萄'),('OpenClaw带火AI记忆',510,'机器之心'),('美军AI：Claude锁定1000目标',385,'新智元'),('Claude设计负责人谈AI',302,'机器之心'),('火山ArkClaw一人指挥4AI',287,'新智元'),('陶哲轩专访AI×Science',113,'量子位')]
mh=arts[0][1];cx2=LM+Inches(4.5);bmw2=Inches(5.5);bh3=Inches(0.35);bg2=Inches(0.12)
hy2=Inches(1.2);at(s,LM,hy2,Inches(0.5),Inches(0.3),'#',fs=Pt(12),fc=NAVY,b=True)
at(s,LM+Inches(0.5),hy2,Inches(3.8),Inches(0.3),'标题',fs=Pt(12),fc=NAVY,b=True)
at(s,cx2,hy2,Inches(3),Inches(0.3),'热度指数',fs=Pt(12),fc=NAVY,b=True)
ahl(s,LM,hy2+Inches(0.3),CW,BLACK,Pt(0.75))
for i,(ti,ht,md) in enumerate(arts):
    by=Inches(1.65)+i*(bh3+bg2);bw3=int(bmw2*ht/mh);co=NAVY if i==0 else(ACCENT_BLUE if i<3 else LINE_GRAY)
    at(s,LM,by,Inches(0.5),bh3,str(i+1),fs=Pt(12),fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE)
    at(s,LM+Inches(0.5),by,Inches(3.8),bh3,ti,fs=Pt(11),fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
    ar(s,cx2,by+Inches(0.03),bw3,bh3-Inches(0.06),co)
    at(s,cx2+bw3+Inches(0.1),by,Inches(1.5),bh3,f'{ht}  ({md})',fs=Pt(9),fc=MED_GRAY,an=MSO_ANCHOR.MIDDLE)
asrc(s,'Source: 科技垂类媒体热文监控，2026-03-13 至 2026-03-15');apn(s,12,TT)

# ═══ Slide 12: 企业决策者 — 模板 #61 Checklist/Status ═══
s=prs.slides.add_slide(BL)
aat(s,'对企业决策者：六大行动项进度追踪——P0级需立即执行',ts=Pt(20))
hy3=Inches(1.3);hds=[('#',Inches(0.5)),('行动项',Inches(4.5)),('优先级',Inches(1.2)),('时间窗口',Inches(1.8)),('状态',Inches(2.0))]
hx3=LM
for lb,w in hds:at(s,hx3,hy3,w,Inches(0.35),lb,fs=Pt(12),fc=NAVY,b=True);hx3+=w
ahl(s,LM,hy3+Inches(0.35),CW,BLACK,Pt(0.75))
tks=[('1','设立实权CAIO，启动Agent化试点','P0','立即','active'),('2','建立M2架构中台','P0','30天内','pending'),('3','构建私有化数据资产库','P1','60天内','pending'),('4','评估315合规风险','P0','1周内','active'),('5','探索OpenClaw生态','P2','持续','pending'),('6','HR重定义人机边界','P1','90天内','pending')]
sc={'done':('✓ 完成',ACCENT_GREEN,LIGHT_GREEN),'active':('→ 进行中',ACCENT_ORANGE,LIGHT_ORANGE),'pending':('○ 待启动',MED_GRAY,BG_GRAY)}
rh3=Inches(0.75)
for i,(nm,tk,pr,dl,st) in enumerate(tks):
    ry=Inches(1.8)+i*rh3;sl,scl,sbg=sc[st]
    if i%2==0:ar(s,LM,ry,CW,rh3,RGBColor(0xFA,0xFA,0xFA))
    at(s,LM,ry,Inches(0.5),rh3,nm,fs=BODY_SIZE,fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
    at(s,LM+Inches(0.5),ry,Inches(4.5),rh3,tk,fs=BODY_SIZE,fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
    pc=ACCENT_RED if pr=='P0'else(ACCENT_ORANGE if pr=='P1'else ACCENT_BLUE)
    ar(s,LM+Inches(5.15),ry+Inches(0.12),Inches(0.8),rh3-Inches(0.24),pc)
    at(s,LM+Inches(5.15),ry,Inches(0.8),rh3,pr,fs=Pt(11),fc=WHITE,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    at(s,LM+Inches(6.2),ry,Inches(1.8),rh3,dl,fs=BODY_SIZE,fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
    ar(s,LM+Inches(8.15),ry+Inches(0.1),Inches(1.5),rh3-Inches(0.2),sbg)
    at(s,LM+Inches(8.15),ry,Inches(1.5),rh3,sl,fs=Pt(12),fc=scl,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    if i<len(tks)-1:ahl(s,LM,ry+rh3,CW,LINE_GRAY,Pt(0.25))
dn=sum(1 for t in tks if t[4]=='done');ac=sum(1 for t in tks if t[4]=='active');pg=(dn+ac*0.5)/len(tks)
py2=Inches(6.35);ar(s,LM,py2,CW,Inches(0.6),BG_GRAY)
at(s,LM+Inches(0.3),py2,Inches(3),Inches(0.6),f'总进度：{dn}/{len(tks)} 完成，{ac} 进行中 ({int(pg*100)}%)',fs=BODY_SIZE,fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE)
bx3=LM+Inches(3.5);bw4=CW-Inches(4);ar(s,bx3,py2+Inches(0.2),bw4,Inches(0.2),LINE_GRAY);ar(s,bx3,py2+Inches(0.2),int(bw4*pg),Inches(0.2),ACCENT_ORANGE)
asrc(s,'Source: 战略分析团队建议，2026年3月15日');apn(s,13,TT)

# ═══ Slide 13: 投资者 — 模板 #30 Vertical Steps ═══
s=prs.slides.add_slide(BL)
aat(s,'对投资者：三阶段布局——短期安全→中期Agent基础设施→长期CaaS平台')
at(s,LM,Inches(1.15),CW,Inches(0.25),'垂直步骤图 — 投资布局时间线',fs=Pt(10),fc=MED_GRAY)
lgx13=LM+Inches(4.0);ar(s,lgx13,Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_RED);at(s,lgx13+Inches(0.2),Inches(1.15),Inches(1.2),Inches(0.25),'短期催化（红）',fs=Pt(10),fc=MED_GRAY)
ar(s,lgx13+Inches(1.6),Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_BLUE);at(s,lgx13+Inches(1.8),Inches(1.15),Inches(1.2),Inches(0.25),'中期主线（蓝）',fs=Pt(10),fc=MED_GRAY)
ar(s,lgx13+Inches(3.2),Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_GREEN);at(s,lgx13+Inches(3.4),Inches(1.15),Inches(1.2),Inches(0.25),'长期赛道（绿）',fs=Pt(10),fc=MED_GRAY)
sis=[('短期催化','1-3个月',ACCENT_RED,LIGHT_RED,'315曝光利好AI安全板块。关注：国投智能、三六零、格灵深瞳'),('中期主线','3-12个月',ACCENT_BLUE,LIGHT_BLUE,'AI Agent基础设施。关注：多智能体编排、AI记忆系统、Agent工具链'),('长期赛道','1-3年',ACCENT_GREEN,LIGHT_GREEN,'企业M2架构服务商、行业垂直智能体平台、CaaS能力软件化平台')]
sys2=Inches(1.65);sth=Inches(1.35);stg=Inches(0.15);tlx=LM+Inches(0.3);tlw=Inches(0.06)
for i,(ti,pd,ac,lbg,de) in enumerate(sis):
    sy2=sys2+i*(sth+stg);ar(s,tlx,sy2,tlw,sth,ac);ao(s,tlx-Inches(0.2),sy2,str(i+1),Inches(0.45),ac)
    cdx=LM+Inches(0.9);cdw=CW-Inches(1.1);ar(s,cdx,sy2,cdw,sth,lbg);ar(s,cdx,sy2,cdw,Inches(0.06),ac)
    at(s,cdx+Inches(0.2),sy2+Inches(0.15),Inches(2.5),Inches(0.35),ti,fs=SUB_HEADER_SIZE,fc=ac,b=True)
    at(s,cdx+Inches(2.8),sy2+Inches(0.15),Inches(1.5),Inches(0.35),pd,fs=Pt(12),fc=MED_GRAY,b=True)
    at(s,cdx+Inches(0.2),sy2+Inches(0.6),cdw-Inches(0.4),Inches(0.6),de,fs=Pt(12),fc=DARK_GRAY)
ar(s,LM,Inches(6.35),CW,Inches(0.55),BG_GRAY)
at(s,LM+Inches(0.3),Inches(6.35),Inches(1.5),Inches(0.55),'风险提示',fs=BODY_SIZE,fc=ACCENT_RED,b=True,an=MSO_ANCHOR.MIDDLE)
at(s,LM+Inches(2),Inches(6.35),CW-Inches(2.3),Inches(0.55),'监管加速风险 | 泡沫信号（讨论焦点仅限科技公司本身即为泡沫）| 人才断层5年后显现',fs=Pt(12),fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
asrc(s,'Source: 综合分析，热点中台事件库');apn(s,14,TT)

# ═══ Slide 14: Agent平台 — 模板 #59 Stakeholder Map ═══
s=prs.slides.add_slide(BL)
aat(s,'深度对比：四大Agent平台竞争矩阵——用户规模 vs 技术壁垒',ts=Pt(20))
at(s,LM,Inches(1.15),CW,Inches(0.25),'2×2矩阵图 — X轴: 技术壁垒 | Y轴: 用户规模',fs=Pt(10),fc=MED_GRAY)
lgx14=LM+Inches(5.2);ar(s,lgx14,Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_RED);at(s,lgx14+Inches(0.2),Inches(1.15),Inches(1.0),Inches(0.25),'核心赢家',fs=Pt(10),fc=MED_GRAY)
ar(s,lgx14+Inches(1.3),Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_GREEN);at(s,lgx14+Inches(1.5),Inches(1.15),Inches(0.8),Inches(0.25),'观望区',fs=Pt(10),fc=MED_GRAY)
ar(s,lgx14+Inches(2.5),Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_BLUE);at(s,lgx14+Inches(2.7),Inches(1.15),Inches(1.2),Inches(0.25),'Keep Informed',fs=Pt(10),fc=MED_GRAY)
ar(s,lgx14+Inches(4.1),Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_ORANGE);at(s,lgx14+Inches(4.3),Inches(1.15),Inches(0.9),Inches(0.25),'潜力选手',fs=Pt(10),fc=MED_GRAY)
gl2=LM+Inches(2.0);gt2=Inches(1.65);cw4=Inches(4.5);ch4=Inches(2.0)
qlb=[('高用户+低壁垒','Keep Informed',LIGHT_BLUE),('高用户+高壁垒','核心赢家',LIGHT_RED),('低用户+低壁垒','观察区',LIGHT_GREEN),('低用户+高壁垒','潜力选手',LIGHT_ORANGE)]
qst=[[('字节 ArkClaw','飞书/抖音生态，热度287')],[('腾讯 QClaw','社交入口垄断，765'),('阿里 JVS','云服务生态，599')],[('其他跟随者','碎片化竞争')],[('OpenClaw','开源框架，160K+Stars')]]
for qi,(lc2,le,bg) in enumerate(qlb):
    ro=qi//2;co=qi%2;qx=gl2+co*cw4;qy=gt2+ro*ch4
    ar(s,qx,qy,cw4-Inches(0.05),ch4-Inches(0.05),bg)
    at(s,qx+Inches(0.15),qy+Inches(0.1),cw4-Inches(0.3),Inches(0.35),lc2,fs=Pt(13),fc=NAVY,b=True)
    for ni,(nm,dt) in enumerate(qst[qi]):
        at(s,qx+Inches(0.15),qy+Inches(0.5)+ni*Inches(0.5),cw4-Inches(0.3),Inches(0.25),nm,fs=Pt(12),fc=DARK_GRAY,b=True)
        at(s,qx+Inches(0.15),qy+Inches(0.75)+ni*Inches(0.5),cw4-Inches(0.3),Inches(0.2),dt,fs=Pt(10),fc=MED_GRAY)
at(s,LM,gt2,Inches(1.8),2*ch4,'用\n户\n规\n模\n↑',fs=Pt(12),fc=MED_GRAY,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
at(s,gl2,gt2+2*ch4+Inches(0.1),2*cw4,Inches(0.3),'技术壁垒 →',fs=Pt(12),fc=MED_GRAY,al=PP_ALIGN.CENTER)
ar(s,LM,Inches(6.0),CW,Inches(0.9),BG_GRAY)
at(s,LM+Inches(0.3),Inches(6.0),CW-Inches(0.6),Inches(0.9),'Agent平台之争本质是"入口之争"。腾讯握社交，阿里握云，字节握内容。OpenClaw可能成为三家通吃的底座。',fs=BODY_SIZE,fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE)
asrc(s,'Source: 知识库监控报告，GitHub，产品公开信息');apn(s,15,TT)

# ═══ Slide 15: 效率陷阱 — 模板 #60 Issue/Decision Tree ═══
s=prs.slides.add_slide(BL)
aat(s,'"效率陷阱"深度剖析：问题诊断树',ts=Pt(20))
at(s,LM,Inches(1.15),CW,Inches(0.25),'决策树 — 从核心矛盾到具体表现的因果分解',fs=Pt(10),fc=MED_GRAY)
lgx15=LM+Inches(5.0);ar(s,lgx15,Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_BLUE);at(s,lgx15+Inches(0.2),Inches(1.15),Inches(1.4),Inches(0.25),'冗余性悖论（蓝）',fs=Pt(10),fc=MED_GRAY)
ar(s,lgx15+Inches(1.8),Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_ORANGE);at(s,lgx15+Inches(2.0),Inches(1.15),Inches(1.3),Inches(0.25),'KPI通胀（橙）',fs=Pt(10),fc=MED_GRAY)
ar(s,lgx15+Inches(3.5),Inches(1.18),Inches(0.15),Inches(0.15),ACCENT_RED);at(s,lgx15+Inches(3.7),Inches(1.15),Inches(1.3),Inches(0.25),'梯队断层（红）',fs=Pt(10),fc=MED_GRAY)
# Helper: draw a proper connector line from point to point
def aln(s,x1,y1,x2,y2,c=LINE_GRAY,th=Pt(1.5)):
    """Draw line from (x1,y1) to (x2,y2) using a freeform/connector approach"""
    w=abs(x2-x1)if x2!=x1 else Emu(12700);h=abs(y2-y1)if y2!=y1 else Emu(12700)
    lx=min(x1,x2);ly=min(y1,y2)
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,lx,ly,w,h)
    ln.fill.background();ln.line.fill.background();_cs(ln)
    # Use thin rectangles for horizontal and vertical segments
    if abs(int(x2)-int(x1))>abs(int(y2)-int(y1)):
        r=ar(s,min(x1,x2),y1,abs(x2-x1),Emu(int(th*1.2)),c)
    else:
        r=ar(s,x1,min(y1,y2),Emu(int(th*1.2)),abs(y2-y1),c)
    return r
# L0: Root node
L0x=LM+Inches(0.2);L0y=Inches(2.8);L0w=Inches(2.0);L0h=Inches(1.2)
ar(s,L0x,L0y,L0w,L0h,NAVY);at(s,L0x,L0y,L0w,L0h,'效率陷阱\n核心矛盾',fs=Pt(16),fc=WHITE,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
# L0 right edge midpoint
L0rx=L0x+L0w;L0my=L0y+L0h//2
# L1: Three sub-issues
L1s=[('冗余性悖论','→自证冗余'),('KPI通胀','→薪资停滞'),('梯队断层','→人才消失')]
bcs=[ACCENT_BLUE,ACCENT_ORANGE,ACCENT_RED]
L1x=L0rx+Inches(0.7);L1w=Inches(1.8);L1h=Inches(0.9)
L1ys=[Inches(1.55),Inches(2.8),Inches(4.05)]
# Vertical trunk line from L0 right edge
trunk_x=L0rx+Inches(0.35)
aln(s,L0rx,L0my,trunk_x,L0my,LINE_GRAY,Pt(1.5))
aln(s,trunk_x,L1ys[0]+L1h//2,trunk_x,L1ys[2]+L1h//2,LINE_GRAY,Pt(1.5))
for i,(ti,mt) in enumerate(L1s):
    L1y=L1ys[i]
    # Horizontal branch from trunk to L1 box left edge
    aln(s,trunk_x,L1y+L1h//2,L1x,L1y+L1h//2,LINE_GRAY,Pt(1.5))
    ar(s,L1x,L1y,L1w,L1h,bcs[i])
    at(s,L1x,L1y,L1w,Inches(0.5),ti,fs=Pt(14),fc=WHITE,b=True,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    at(s,L1x,L1y+Inches(0.48),L1w,Inches(0.4),mt,fs=Pt(11),fc=WHITE,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
# L2: Leaf nodes (2 per L1)
L2g=[[('展示AI价值','→被裁'),('产出3倍','→砍2/3')],[('KPI成倍增','→薪资不变'),('按Token计','→白领困境')],[('初级岗消失','→学徒崩'),('5年后断层','→系统风险')]]
L2x=L1x+L1w+Inches(0.7);L2w=Inches(1.6);L2h=Inches(0.55)
for gi,gp in enumerate(L2g):
    L1y_center=L1ys[gi]+L1h//2;L1rx=L1x+L1w
    # Trunk from L1 right edge
    t2x=L1rx+Inches(0.35)
    aln(s,L1rx,L1y_center,t2x,L1y_center,LINE_GRAY,Pt(1))
    leaf_ys=[L1ys[gi]-Inches(0.05),L1ys[gi]+L1h-L2h+Inches(0.05)]
    if len(gp)>1:aln(s,t2x,leaf_ys[0]+L2h//2,t2x,leaf_ys[1]+L2h//2,LINE_GRAY,Pt(1))
    for li,(ti,mt) in enumerate(gp):
        L2y=leaf_ys[li]
        aln(s,t2x,L2y+L2h//2,L2x,L2y+L2h//2,LINE_GRAY,Pt(1))
        ar(s,L2x,L2y,L2w,L2h,BG_GRAY)
        at(s,L2x+Inches(0.08),L2y,L2w*0.55,L2h,ti,fs=Pt(10),fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
        at(s,L2x+L2w*0.55,L2y,L2w*0.45,L2h,mt,fs=Pt(10),fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE,al=PP_ALIGN.CENTER)
# Right panel: response strategies
L3x=L2x+L2w+Inches(0.25);L3w=CW-(L3x-LM)
ar(s,L3x,Inches(1.5),L3w,Inches(4.2),BG_GRAY)
at(s,L3x+Inches(0.15),Inches(1.6),L3w-Inches(0.3),Inches(0.35),'应对策略',fs=Pt(14),fc=NAVY,b=True)
ahl(s,L3x+Inches(0.15),Inches(2.0),L3w-Inches(0.3),LINE_GRAY)
at(s,L3x+Inches(0.15),Inches(2.15),L3w-Inches(0.3),Inches(3.2),['1. 重新定义人机协作边界','','2. AI作为增强而非替代','','3. 保留初级培训岗位','','4. KPI与AI工具脱钩','','5. 建立人才储备基金'],fs=Pt(12),fc=DARK_GRAY,ls=Pt(3))
ar(s,LM,Inches(6.0),CW,Inches(0.9),ACCENT_RED)
at(s,LM+Inches(0.3),Inches(6.0),CW-Inches(0.6),Inches(0.9),'战略警示：HR必须在90天内重新定义人机职责边界，否则面临"人才断层+劳动法风险"双重困境。',fs=BODY_SIZE,fc=WHITE,b=True,an=MSO_ANCHOR.MIDDLE)
asrc(s,'Source: PwC 2025, Nimdzi 2025');apn(s,16,TT)

# ═══ Slide 16: Key Takeaway — 模板 #25 Key Takeaway with Detail ═══
s=prs.slides.add_slide(BL)
aat(s,'核心洞见：行业正从"大模型崇拜"全面转入"系统重构"深水区')
lwk=Inches(7.5)
at(s,LM,Inches(1.4),lwk,Inches(0.4),'三条主线同步演化',fs=SUB_HEADER_SIZE,fc=NAVY,b=True)
ahl(s,LM,Inches(1.9),lwk,LINE_GRAY)
at(s,LM,Inches(2.1),lwk,Inches(4),['执行时代开启：','AI Agent升级为"桌面操作系统"。OpenClaw 16万星标，腾讯/阿里/字节入场终局。','','架构决定成败：','95%AI项目败于架构。"双机理论"M2架构融合层成为核心壁垒。','','安全与治理：','315晚会首次将AI安全作为核心议题。监管从观望转向主动介入。'],fs=Pt(12),fc=DARK_GRAY,ls=Pt(3))
tkx=LM+lwk+Inches(0.3);tkw=CW-lwk-Inches(0.3)
ar(s,tkx,Inches(1.4),tkw,Inches(5.2),BG_GRAY)
at(s,tkx+Inches(0.2),Inches(1.6),tkw-Inches(0.4),Inches(0.4),'Key Takeaways',fs=BODY_SIZE,fc=NAVY,b=True)
ahl(s,tkx+Inches(0.2),Inches(2.1),tkw-Inches(0.4),LINE_GRAY)
tks2=[('1','Agent时代已至','对话→执行不可逆',ACCENT_BLUE),('2','架构>模型','M2层决定AI生死',ACCENT_GREEN),('3','安全强监管','315标志新阶段',ACCENT_ORANGE),('4','组织重塑','CaaS元年开启',ACCENT_RED),('5','制造业30%+','AI成新常态',NAVY)]
for i,(nm,ti,de,cl) in enumerate(tks2):
    ty2=Inches(2.3)+i*Inches(0.85);ao(s,tkx+Inches(0.2),ty2,nm,Inches(0.35),cl)
    at(s,tkx+Inches(0.65),ty2,tkw-Inches(0.85),Inches(0.3),ti,fs=Pt(13),fc=cl,b=True)
    at(s,tkx+Inches(0.65),ty2+Inches(0.32),tkw-Inches(0.85),Inches(0.35),de,fs=Pt(11),fc=MED_GRAY)
asrc(s,'Source: 综合分析，知识库监控报告');apn(s,17,TT)

# ═══ Slide 17: 方法论 — 模板 #58 Dashboard Table + Factoids ═══
s=prs.slides.add_slide(BL)
aat(s,'方法论说明：覆盖9大平台、2,571+条信号、McKinsey MECE分析框架',ts=Pt(20))
lwm=Inches(6.2);tym=Inches(1.2);cns=['维度','说明'];cws=[Inches(1.8),Inches(4)]
hxm=LM
for cn,cw5 in zip(cns,cws):at(s,hxm,tym,cw5,Inches(0.3),cn,fs=Pt(12),fc=NAVY,b=True);hxm+=cw5
ahl(s,LM,tym+Inches(0.3),lwm,BLACK,Pt(0.75))
mds=[('数据采集','9大平台自动监控'),('采集频率','每12小时，3日6轮次'),('分析框架','MECE + FBTA四层法'),('热度计算','传播量×速度×权重'),('舆情模型','正面/中性/负面三维')]
for ri,(dm,de) in enumerate(mds):
    rym=tym+Inches(0.4)+ri*Inches(0.5)
    at(s,LM,rym,cws[0],Inches(0.4),dm,fs=BODY_SIZE,fc=NAVY,b=True,an=MSO_ANCHOR.MIDDLE)
    at(s,LM+cws[0],rym,cws[1],Inches(0.4),de,fs=BODY_SIZE,fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
    if ri<len(mds)-1:ahl(s,LM,rym+Inches(0.45),lwm,LINE_GRAY,Pt(0.25))
# Right: mini horizontal bar chart
chx=LM+lwm+Inches(0.5);chw=CW-lwm-Inches(0.5);chy=Inches(1.3)
at(s,chx,chy,chw,Inches(0.3),'各平台信号量占比',fs=Pt(12),fc=NAVY,b=True)
ahl(s,chx,chy+Inches(0.3),chw,BLACK,Pt(0.5))
plts=[('微博',420),('B站',380),('抖音',350),('快手',280),('头条',250),('小红书',230),('Reddit',240),('Twitter',220),('YouTube',201)]
mxp=420;bm3=Inches(2.8);bh4=Inches(0.22);bg3=Inches(0.08)
for i,(nm,vl) in enumerate(plts):
    by3=chy+Inches(0.45)+i*(bh4+bg3);bw5=int(bm3*vl/mxp)
    at(s,chx,by3,Inches(0.8),bh4,nm,fs=Pt(9),fc=MED_GRAY,an=MSO_ANCHOR.MIDDLE)
    ar(s,chx+Inches(0.9),by3,bw5,bh4,NAVY if i<3 else ACCENT_BLUE)
    at(s,chx+Inches(0.9)+bw5+Inches(0.05),by3,Inches(0.5),bh4,str(vl),fs=Pt(8),fc=DARK_GRAY,an=MSO_ANCHOR.MIDDLE)
# Factoid cards (#58 bottom)
fts=[('6份','知识库报告',ACCENT_BLUE),('177条','热点事件',NAVY),('30篇','科技热文',ACCENT_GREEN),('9大','覆盖平台',ACCENT_ORANGE)]
fw=CW/len(fts)-Inches(0.15);fh=Inches(1.3);fy=Inches(5.5)
for i,(vl,lb,cl) in enumerate(fts):
    fx=LM+i*(fw+Inches(0.15));ar(s,fx,fy,fw,fh,BG_GRAY);ar(s,fx,fy,Inches(0.06),fh,cl)
    at(s,fx+Inches(0.2),fy+Inches(0.1),fw-Inches(0.3),Inches(0.5),vl,fs=Pt(24),fc=cl,b=True)
    at(s,fx+Inches(0.2),fy+Inches(0.6),fw-Inches(0.3),Inches(0.3),lb,fs=Pt(11),fc=MED_GRAY)
asrc(s,'数据截止：2026年3月15日 21:57 CST');apn(s,18,TT)

# ═══ Slide 18: Closing ═══
s=prs.slides.add_slide(BL)
ar(s,0,0,SW,Inches(0.05),NAVY)
at(s,Inches(1.5),Inches(2.0),Inches(10.3),Inches(1),'Agent接管，架构制胜，安全筑基',fs=Pt(28),fc=NAVY,b=True,fn='Georgia',al=PP_ALIGN.CENTER)
ahl(s,Inches(5.5),Inches(3.3),Inches(2.3),NAVY,Pt(1.5))
at(s,Inches(1.5),Inches(3.8),Inches(10.3),Inches(2),'从大模型崇拜到系统重构\nAgent时代已至，架构决定成败\n安全治理从边缘走向中心',fs=SUB_HEADER_SIZE,fc=DARK_GRAY,al=PP_ALIGN.CENTER,ls=Pt(14))
at(s,Inches(1.5),Inches(6.0),Inches(10.3),Inches(0.4),'CONFIDENTIAL  |  AI行业三日趋势深度报告  |  2026年3月',fs=Pt(12),fc=MED_GRAY,al=PP_ALIGN.CENTER)
ahl(s,Inches(1),Inches(6.8),Inches(3),NAVY,Pt(2));apn(s,19,TT)

# ═══ Save ═══
op=os.path.join(os.path.dirname(os.path.abspath(__file__)),'output','AI行业三日趋势深度报告_20260313-15.pptx')
os.makedirs(os.path.dirname(op),exist_ok=True);prs.save(op);fc(op)
print(f'✅ Saved: {op}')
